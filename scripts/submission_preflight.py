"""Read-only preflight for the 2026 금융 AI Challenge handoff.

Without ``--strict`` the command distinguishes operator-owned pending work
(participant identity and public URL) from repository failures. With
``--strict --url ...`` it is the final go/no-go check immediately before the
operator uploads the entries to Daker.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import ssl
import sys
import urllib.error
import urllib.parse
import urllib.request
import zipfile
from dataclasses import dataclass
from pathlib import Path

REPO_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(REPO_ROOT / "scripts"))
from _console import force_utf8_stdout  # noqa: E402

force_utf8_stdout()

COMPETITION_DIR = REPO_ROOT / "docs" / "competition"
OUTPUT_DIR = REPO_ROOT / "output"
LOCAL_PROFILE = COMPETITION_DIR / "submission_profile.local.json"
EXAMPLE_PROFILE = COMPETITION_DIR / "submission_profile.example.json"
PLANNING_PDF = OUTPUT_DIR / "pdf" / "2026_금융_AI_Challenge_기획서_Home_Compass.pdf"
FEATURE_PDF = OUTPUT_DIR / "pdf" / "2026_금융_AI_Challenge_기능명세서_Home_Compass.pdf"
DECK = COMPETITION_DIR / "기술설명서_Home_Compass.pptx"
SOURCE_ZIP = OUTPUT_DIR / "submission" / "Home_Compass_source.zip"
#: 운영자가 채워야 하는 자리는 전부 이 접두어로 시작한다. **정확한 문자열이 아니라
#: 접두어로 거른다** — 자리가 늘 때마다 이 검사를 고쳐야 하면 다음 자리는 안 걸린다.
PLACEHOLDER_PREFIX = "__운영자_"


@dataclass(frozen=True)
class Result:
    status: str
    label: str
    detail: str


def passed(label: str, detail: str) -> Result:
    return Result("PASS", label, detail)


def failed(label: str, detail: str) -> Result:
    return Result("FAIL", label, detail)


def pending(label: str, detail: str) -> Result:
    return Result("PENDING", label, detail)


def check_required_files() -> Result:
    required = [
        REPO_ROOT / "Dockerfile",
        REPO_ROOT / "render.yaml",
        REPO_ROOT / "scripts" / "start_server.py",
        COMPETITION_DIR / "SUBMISSION_RUNBOOK.md",
        COMPETITION_DIR / "build_submission_pdfs.py",
        COMPETITION_DIR / "capture_evidence.py",
        PLANNING_PDF,
        FEATURE_PDF,
        DECK,
        SOURCE_ZIP,
    ]
    missing = [str(path.relative_to(REPO_ROOT)) for path in required if not path.is_file()]
    return failed("required files", ", ".join(missing)) if missing else passed(
        "required files", f"{len(required)} files present"
    )


def check_profile(path: Path) -> Result:
    if not path.is_file():
        return pending("participant identity", f"create {path.relative_to(REPO_ROOT)}")
    try:
        profile = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        return failed("participant identity", f"invalid profile JSON: {exc}")
    team = str(profile.get("team_name") or "").strip()
    members = str(profile.get("member_names") or "").strip()
    if not team or not members or PLACEHOLDER_PREFIX in team + members:
        return pending("participant identity", "registered team/member names are not filled")
    return passed("participant identity", f"team={team}, members={members}")


def check_reviewer_accounts(path: Path) -> Result:
    """심사 계정을 제공하기로 했으면 그 안내가 실제로 채워져 있어야 한다.

    기능명세서는 F7-F10 을 「완료·심사 계정 제공」으로 선언한다. 안내가 비었거나
    자리표시자면 심사위원이 그 넷을 재현할 수 없는데 표에는 완료로 적힌다 —
    **제출물이 스스로 거짓이 되는 자리**다.
    """
    if not path.is_file():
        return pending("reviewer accounts", f"create {path.relative_to(REPO_ROOT)}")
    try:
        profile = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        return failed("reviewer accounts", f"invalid profile JSON: {exc}")
    provided = profile.get("reviewer_accounts_provided")
    if not provided:
        return passed("reviewer accounts", "not provided — F7-F10 are outside the review scope")
    guidance = str(profile.get("reviewer_account_instructions") or "").strip()
    if not guidance or PLACEHOLDER_PREFIX in guidance:
        return pending(
            "reviewer accounts",
            "reviewer_accounts_provided is true but reviewer_account_instructions is empty",
        )
    return passed("reviewer accounts", f"guidance present ({len(guidance)} chars)")


def _pdf_text(path: Path) -> tuple[int, str]:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("install docs/competition/requirements.txt") from exc
    reader = PdfReader(str(path))
    text = "\n".join(page.extract_text() or "" for page in reader.pages)
    return len(reader.pages), text


def check_pdf(path: Path, label: str, headings: list[str], minimum_pages: int) -> list[Result]:
    if not path.is_file():
        return [failed(label, f"missing {path.relative_to(REPO_ROOT)}")]
    try:
        pages, text = _pdf_text(path)
    except Exception as exc:  # artifact corruption and missing verifier both fail closed
        return [failed(label, str(exc))]
    missing = [heading for heading in headings if heading not in text]
    results = []
    if pages < minimum_pages or missing:
        results.append(failed(label, f"pages={pages}, missing headings={missing}"))
    else:
        results.append(passed(label, f"pages={pages}, all official headings found"))
    if PLACEHOLDER_PREFIX in text:
        results.append(pending(f"{label} identity", "PDF still contains the operator placeholder"))
    else:
        results.append(passed(f"{label} identity", "no operator placeholder"))
    return results


def check_deck() -> Result:
    if not DECK.is_file():
        return failed("presentation deck", "PPTX missing")
    try:
        from pptx import Presentation
        deck = Presentation(str(DECK))
        text = "\n".join(
            shape.text for slide in deck.slides for shape in slide.shapes
            if hasattr(shape, "text")
        )
    except Exception as exc:
        return failed("presentation deck", str(exc))
    stale = [phrase for phrase in ("KB 사업 연계", "KB국민은행", "KB 주택금융") if phrase in text]
    if len(deck.slides) != 19 or stale:
        return failed("presentation deck", f"slides={len(deck.slides)}, stale={stale}")
    return passed("presentation deck", "19 slides, previous-contest brand removed")


def check_deployment_config() -> Result:
    docker = (REPO_ROOT / "Dockerfile").read_text(encoding="utf-8")
    render = (REPO_ROOT / "render.yaml").read_text(encoding="utf-8")
    startup = (REPO_ROOT / "scripts" / "start_server.py").read_text(encoding="utf-8")
    required = [
        ("plan: 0.5c-512mb", render),
        ('autoDeployTrigger: "off"', render),
        ("healthCheckPath: /api/health", render),
        ("mountPath: /var/data", render),
        ("sizeGB: 1", render),
        ("numInstances: 1", render),
        ("HOME_COMPASS_COOKIE_SECURE", render),
        ("sync: false", render),
        ("HEALTHCHECK", docker),
        ('"--workers",\n        "1"', startup),
    ]
    missing = [needle for needle, haystack in required if needle not in haystack]
    return failed("deployment config", f"missing invariants: {missing}") if missing else passed(
        "deployment config", "HTTPS cookie, secrets, health check, one instance, persistent disk"
    )


def check_source_zip() -> Result:
    if not SOURCE_ZIP.is_file():
        return failed("source archive", f"missing {SOURCE_ZIP.relative_to(REPO_ROOT)}")
    try:
        with zipfile.ZipFile(SOURCE_ZIP) as archive:
            bad_member = archive.testzip()
            names = archive.namelist()
            manifest_name = "Home_Compass/SUBMISSION_MANIFEST.json"
            manifest = json.loads(archive.read(manifest_name))
            manifest_errors = []
            for entry in manifest.get("files", []):
                archived = "Home_Compass/" + str(entry.get("path") or "")
                if archived not in names:
                    manifest_errors.append(f"missing:{archived}")
                    continue
                actual = hashlib.sha256(archive.read(archived)).hexdigest()
                if actual != entry.get("sha256"):
                    manifest_errors.append(f"sha256:{archived}")
    except (OSError, zipfile.BadZipFile) as exc:
        return failed("source archive", str(exc))
    except (KeyError, json.JSONDecodeError, TypeError) as exc:
        return failed("source archive", f"invalid manifest: {exc}")
    forbidden = [
        name for name in names
        if name.endswith("/.env")
        or name.endswith("submission_profile.local.json")
        or "/.git/" in name
        or "/.venv/" in name
        or "/output/" in name
        or "/tmp/" in name
    ]
    expected_count = manifest.get("fileCount")
    if bad_member or forbidden or manifest_errors or expected_count != len(manifest.get("files", [])):
        return failed(
            "source archive",
            f"bad_member={bad_member}, forbidden={forbidden[:5]}, manifest={manifest_errors[:5]}",
        )
    return passed(
        "source archive",
        # ★ `package_submission.py` 는 **224**, 여기는 **225** 를 말한다. 어긋난 것이
        #   아니라 세는 대상이 다르다 — 매니페스트는 자기 해시를 자기 안에 적을 수 없어
        #   **자기 자신을 목록에서 뺀다.** 두 수를 나란히 보고 사고로 오해하지 않도록
        #   여기서 분모를 밝힌다.
        f"{len(names)} zip entries ({len(names) - 1} files + manifest), "
        "CRC and manifest SHA-256 OK, no local secrets/output",
    )


def _json_request(url: str, payload: dict | None = None) -> tuple[int, object]:
    body = None if payload is None else json.dumps(payload).encode("utf-8")
    request = urllib.request.Request(
        url,
        data=body,
        headers={
            "Accept": "application/json",
            "Content-Type": "application/json",
            "User-Agent": "Home-Compass-Submission-Preflight/1.0",
        },
        method="GET" if payload is None else "POST",
    )
    with urllib.request.urlopen(request, timeout=20, context=ssl.create_default_context()) as response:
        return response.status, json.load(response)


def _text_request(url: str) -> tuple[int, str]:
    request = urllib.request.Request(
        url,
        headers={"Accept": "text/html", "User-Agent": "Home-Compass-Submission-Preflight/1.0"},
    )
    with urllib.request.urlopen(request, timeout=20, context=ssl.create_default_context()) as response:
        return response.status, response.read().decode("utf-8", errors="replace")


def check_live_llm(base: str, sample: dict) -> Result:
    """배포된 서비스가 **실제로** 라이브 LLM 으로 답하는지 끝까지 확인한다.

    `/api/health` 의 `llm` 은 「키가 있고 SDK 를 import 할 수 있다」까지만 말한다.
    모델 id 가 계정에서 거부되면 그 값은 여전히 `openai` 인데 **실제 호출은 실패하고
    응답이 조용히 템플릿으로 내려간다.** 그래서 채팅을 한 번 실제로 불러 `mode` 를 본다 —
    `live` 는 프로바이더가 답했다는 뜻이고 `offline` 은 템플릿이 답했다는 뜻이다.

    라이브를 쓰지 않기로 했다면 이 항목은 `PENDING` 으로 남고, 그것이 맞는 상태다.
    """
    try:
        status, body = _json_request(
            base + "/api/chat",
            {"message": "월세와 전세 중 무엇이 유리한가요?", "profile": sample, "history": []},
        )
    except (OSError, ValueError, urllib.error.URLError) as exc:
        return failed("live LLM", f"/api/chat request failed: {exc}")
    if status != 200 or not isinstance(body, dict):
        return failed("live LLM", f"status={status}, body={body!r}")
    mode = str(body.get("mode") or "")
    provider = str(body.get("provider") or "")
    if mode == "live":
        return passed("live LLM", f"chat answered live (provider={provider})")
    return pending(
        "live LLM",
        f"deployed service answers in template mode (mode={mode!r}, provider={provider!r}) — "
        "the API key is missing, rejected, or the model id is not available to this account",
    )


def check_public_url(raw_url: str | None) -> list[Result]:
    if not raw_url:
        return [pending("public URL", "deploy and pass --url https://... before submission")]
    base = raw_url.strip().rstrip("/")
    parsed = urllib.parse.urlparse(base)
    if parsed.scheme != "https" or not parsed.netloc:
        return [failed("public URL", "an externally submitted URL must be absolute HTTPS")]
    sample = {
        "age": 28,
        "annualIncomeKRW": 42_000_000,
        "monthlyNetIncomeKRW": 3_000_000,
        "liquidAssetsKRW": 40_000_000,
        "existingDebtMonthlyKRW": 300_000,
        "householdSize": 1,
        "regionCode": "11440",
        "isHomeless": True,
        "isNewlywed": False,
        "isSMEEmployee": True,
        "preferredType": "any",
    }
    try:
        home_status, home = _text_request(base + "/")
        health_status, health = _json_request(base + "/api/health")
        analyze_status, result = _json_request(base + "/api/analyze", sample)
    except (OSError, ValueError, urllib.error.URLError) as exc:
        return [failed("public URL", f"request failed: {exc}")]
    if home_status != 200 or "Home_Compass" not in home:
        return [failed("public URL homepage", f"status={home_status}, expected Home_Compass HTML")]
    if health_status != 200 or not isinstance(health, dict) or health.get("status") != "ok":
        return [failed("public URL health", f"status={health_status}, body={health!r}")]
    affordability = result.get("affordability", {}) if isinstance(result, dict) else {}
    expected = (860_000, 730_000)
    actual = (
        affordability.get("maxMonthlyHousingCostKRW"),
        affordability.get("recommendedMonthlyHousingCostKRW"),
    )
    if analyze_status != 200 or actual != expected:
        return [failed("public URL analysis", f"status={analyze_status}, expected={expected}, actual={actual}")]
    required_keys = {"scenarios", "policies", "risk", "provenance", "dataGrade"}
    missing = sorted(required_keys - set(result))
    if missing:
        return [failed("public URL analysis", f"missing response keys: {missing}")]
    return [
        passed("public URL homepage", f"{base}/ -> Home_Compass HTML"),
        passed("public URL health", f"{base}/api/health -> ok"),
        passed("public URL analysis", "sample profile -> 860,000 / 730,000 KRW and complete evidence"),
        # ★ 마지막에 부른다. 위 셋이 깨지면 라이브 여부를 물을 것도 없고, 그때는 호출을
        #   아끼는 편이 낫다 — 이 한 줄이 실제 LLM 요금을 쓰는 유일한 검사다.
        check_live_llm(base, sample),
    ]


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Check Home_Compass submission readiness")
    parser.add_argument("--strict", action="store_true", help="treat operator-owned pending items as failures")
    parser.add_argument("--url", help="deployed public HTTPS origin")
    parser.add_argument("--profile", type=Path, help="participant profile JSON")
    args = parser.parse_args(argv)

    profile = args.profile.resolve() if args.profile else (
        LOCAL_PROFILE if LOCAL_PROFILE.is_file() else EXAMPLE_PROFILE
    )
    results = [
        check_required_files(),
        check_profile(profile),
        check_reviewer_accounts(profile),
    ]
    results += check_pdf(
        PLANNING_PDF,
        "planning PDF",
        ["서비스 명칭", "아이디어 기획 핵심내용", "문제 정의 및 제안 배경", "서비스 컨셉 및 차별성", "활용 데이터 및 생성형 AI 모델 적용 방안", "기대 효과 및 확장 가능성"],
        3,
    )
    results += check_pdf(
        FEATURE_PDF,
        "feature PDF",
        ["MVP 구현 범위", "주요 기능 목록", "사용자 이용 흐름", "AI 및 데이터 처리 방식", "MVP 검증 방법"],
        4,
    )
    results += [check_deck(), check_deployment_config(), check_source_zip()]
    results += check_public_url(args.url)

    for result in results:
        print(f"[{result.status:7}] {result.label}: {result.detail}")
    failures = [result for result in results if result.status == "FAIL"]
    pendings = [result for result in results if result.status == "PENDING"]
    print(f"SUMMARY pass={sum(r.status == 'PASS' for r in results)} pending={len(pendings)} fail={len(failures)}")
    if failures or (args.strict and pendings):
        return 1
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
