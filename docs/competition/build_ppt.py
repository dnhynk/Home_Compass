# -*- coding: utf-8 -*-
"""
build_ppt.py - Technical specification deck builder for "Home_Compass".

Generates docs/competition/기술설명서_Home_Compass.pptx (16:9, 19 slides, 맑은 고딕).
All architecture / engine-flow diagrams are drawn with real pptx shapes
(rectangles, arrows, connectors) - never as plain bullet text.

Usage:
    python build_ppt.py            # build + self-verify
    python build_ppt.py --verify   # verify an already-built file only
"""

from __future__ import annotations

import json
import math
import os
import sys

# 출력 인코딩 가드. 이 스크립트는 한국어를 그대로 찍으므로 출력을 파이프·파일로 넘기는
# 순간 cp949 에서 죽는다 — 실제로 `--verify` 가 그렇게 죽고 있었다. `scripts/_console.py`
# 가 그 규약의 구현이고 **여기서 다시 만들지 않는다** (같은 규약이 세 벌이 되면 한쪽만
# 고쳐진다). `scripts/` 는 패키지가 아니라 최상위 모듈 자리라 경로를 직접 얹는다.
_REPO_ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
sys.path.insert(0, os.path.join(_REPO_ROOT, "scripts"))
from _console import force_utf8_stdout  # noqa: E402

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Emu, Inches, Pt

# --------------------------------------------------------------------------
# constants
# --------------------------------------------------------------------------

HERE = os.path.dirname(os.path.abspath(__file__))
OUT_PATH = os.path.join(HERE, "기술설명서_Home_Compass.pptx")

FONT = "맑은 고딕"

SLIDE_W = 13.333
SLIDE_H = 7.5

# Home_Compass visual system: warm yellow + dark grey.
YELLOW = RGBColor(0xFF, 0xB8, 0x00)
YELLOW_DK = RGBColor(0xD9, 0x9A, 0x00)
YELLOW_LT = RGBColor(0xFF, 0xF4, 0xD6)
DARK = RGBColor(0x2B, 0x2B, 0x2B)
DARK_2 = RGBColor(0x44, 0x44, 0x44)
GRAY = RGBColor(0x70, 0x74, 0x7A)
GRAY_LT = RGBColor(0x9A, 0x9E, 0xA5)
LINE = RGBColor(0xDC, 0xDF, 0xE4)
PANEL = RGBColor(0xF5, 0xF6, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NAVY = RGBColor(0x1F, 0x35, 0x54)
NAVY_LT = RGBColor(0xE8, 0xED, 0xF3)
GREEN = RGBColor(0x2E, 0x7D, 0x5B)
GREEN_LT = RGBColor(0xE4, 0xF1, 0xEB)
AMBER = RGBColor(0xC9, 0x8A, 0x00)
RED = RGBColor(0xC0, 0x39, 0x2B)
RED_LT = RGBColor(0xFA, 0xE9, 0xE7)

CONTENT_X = 0.65
CONTENT_W = 12.03  # 0.65 .. 12.68
BODY_TOP = 1.52
FOOTER_Y = 6.92

#: 슬라이드 수. **손으로 세지 않는다** — `build()` 의 목록 길이가 정본이고 여기는 그것을
#: 받아 쓴다. 15장이던 시절 이 수가 여러 자리에 박혀 있었고, 본선 산출물 넷을 더하면서
#: 그 자리들이 전부 낡을 뻔했다 (이 덱이 「15종」에서 이미 겪은 부류다).
TOTAL_SLIDES = 0  # build() 가 채운다
FOOTER_TEXT = "Home_Compass  ·  2026 금융 AI Challenge 기술설명서"

DISCLAIMER = (
    "프로토타입 시연용 예시 수치입니다. 실제 조건은 취급 금융기관 고시 기준을 따릅니다."
)


# --------------------------------------------------------------------------
# low-level helpers
# --------------------------------------------------------------------------


def _set_typeface(run, name: str) -> None:
    """Force latin/eastasian/complex-script typeface so Hangul renders in 맑은 고딕."""
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(qn("a:latin"))
    if latin is None:
        latin = parse_xml('<a:latin %s/>' % nsdecls("a"))
        rPr.append(latin)
    latin.set("typeface", name)
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = parse_xml("<%s %s/>" % (tag, nsdecls("a")))
            latin.addnext(el) if tag == "a:ea" else rPr.find(qn("a:ea")).addnext(el)
        el.set("typeface", name)


def style_run(run, size=12, bold=False, color=DARK, italic=False) -> None:
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = FONT
    _set_typeface(run, FONT)


def put(tf, items, align=PP_ALIGN.LEFT, line_spacing=1.15):
    """Fill a text frame.

    items: iterable of (text, size, bold, color[, space_before_pt[, align]])
    """
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    first = True
    for it in items:
        text = it[0]
        size = it[1]
        bold = it[2]
        color = it[3]
        space_before = it[4] if len(it) > 4 else 0
        para_align = it[5] if len(it) > 5 else align
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = para_align
        p.space_before = Pt(space_before)
        p.space_after = Pt(0)
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        style_run(run, size, bold, color)
    return tf


def set_insets(tf, left=0.14, top=0.10, right=0.14, bottom=0.10):
    tf.margin_left = Inches(left)
    tf.margin_right = Inches(right)
    tf.margin_top = Inches(top)
    tf.margin_bottom = Inches(bottom)


def txbox(slide, x, y, w, h, items=None, align=PP_ALIGN.LEFT,
          anchor=MSO_ANCHOR.TOP, line_spacing=1.15, insets=(0.0, 0.0, 0.0, 0.0)):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    set_insets(tf, *insets)
    tf.vertical_anchor = anchor
    if items:
        put(tf, items, align=align, line_spacing=line_spacing)
    return box


def rect(slide, x, y, w, h, fill=WHITE, line=LINE, line_w=1.0,
         shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.07,
         items=None, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
         insets=(0.14, 0.10, 0.14, 0.10), line_spacing=1.15, dash=False):
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            shp.adjustments[0] = radius
        except Exception:  # pragma: no cover - defensive
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
        if dash:
            shp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    shp.shadow.inherit = False
    tf = shp.text_frame
    set_insets(tf, *insets)
    tf.vertical_anchor = anchor
    tf.word_wrap = True
    if items:
        put(tf, items, align=align, line_spacing=line_spacing)
    else:
        style_run(tf.paragraphs[0].add_run(), 1, False, fill or WHITE)
    return shp


def arrow(slide, x, y, w, h, direction="right", fill=GRAY_LT):
    shape_type = {
        "right": MSO_SHAPE.RIGHT_ARROW,
        "down": MSO_SHAPE.DOWN_ARROW,
        "left": MSO_SHAPE.LEFT_ARROW,
        "up": MSO_SHAPE.UP_ARROW,
    }[direction]
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    style_run(shp.text_frame.paragraphs[0].add_run(), 1, False, fill)
    return shp


def conn(slide, x1, y1, x2, y2, color=GRAY_LT, width=1.5, arrowhead=True, dash=False):
    c = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if dash:
        c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if arrowhead:
        ln = c.line._get_or_add_ln()
        ln.append(parse_xml('<a:tailEnd %s type="triangle" w="med" len="med"/>' % nsdecls("a")))
    return c


def badge(slide, x, y, w, h, text, fill=YELLOW, color=DARK, size=10.5, bold=True,
          radius=0.5):
    return rect(slide, x, y, w, h, fill=fill, line=None, radius=radius,
                items=[(text, size, bold, color)], insets=(0.05, 0.02, 0.05, 0.02))


# --------------------------------------------------------------------------
# page furniture
# --------------------------------------------------------------------------


def header(slide, title, sub=None, eyebrow=None):
    rect(slide, CONTENT_X, 0.50, 0.10, 0.50, fill=YELLOW, line=None,
         shape_type=MSO_SHAPE.RECTANGLE)
    y = 0.42
    if eyebrow:
        txbox(slide, 0.88, 0.30, 8.0, 0.26,
              [(eyebrow, 10, True, YELLOW_DK)])
        y = 0.55
    txbox(slide, 0.86, y, 11.2, 0.52, [(title, 25, True, DARK)])
    if sub:
        txbox(slide, 0.88, y + 0.56, 11.6, 0.42, [(sub, 12, False, GRAY)])


def footer(slide, idx):
    rect(slide, CONTENT_X, FOOTER_Y, CONTENT_W, 0.014, fill=LINE, line=None,
         shape_type=MSO_SHAPE.RECTANGLE)
    txbox(slide, CONTENT_X, FOOTER_Y + 0.06, 9.5, 0.3, [(FOOTER_TEXT, 8.5, False, GRAY_LT)])
    txbox(slide, 10.9, FOOTER_Y + 0.06, 1.78, 0.3,
          [("%02d / %d" % (idx, TOTAL_SLIDES), 8.5, True, GRAY_LT)], align=PP_ALIGN.RIGHT)


def new_slide(prs, idx, title=None, sub=None, eyebrow=None, chrome=True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if chrome:
        if title:
            header(slide, title, sub, eyebrow)
        footer(slide, idx)
    return slide


def card(slide, x, y, w, h, title, lines, num=None, accent=YELLOW,
         fill=PANEL, title_size=13.5, body_size=11, gap=5):
    box = rect(slide, x, y, w, h, fill=fill, line=LINE, radius=0.06,
               anchor=MSO_ANCHOR.TOP, insets=(0.30, 0.20, 0.22, 0.16), align=PP_ALIGN.LEFT)
    rect(slide, x, y + 0.16, 0.075, h - 0.32, fill=accent, line=None,
         shape_type=MSO_SHAPE.RECTANGLE)
    items = []
    if num:
        items.append((num, 9.5, True, YELLOW_DK))
        items.append((title, title_size, True, DARK, 2))
    else:
        items.append((title, title_size, True, DARK))
    for i, ln in enumerate(lines):
        items.append((ln, body_size, False, DARK_2, gap if i == 0 else 3))
    put(box.text_frame, items, align=PP_ALIGN.LEFT, line_spacing=1.22)
    return box


def note_bar(slide, y, text, h=0.52, fill=YELLOW_LT, color=DARK_2, size=10.5,
             x=CONTENT_X, w=CONTENT_W):
    return rect(slide, x, y, w, h, fill=fill, line=None, radius=0.10,
                items=[(text, size, False, color)], align=PP_ALIGN.LEFT,
                insets=(0.22, 0.08, 0.22, 0.08))


# ==========================================================================
# slides
# ==========================================================================


def slide_01_cover(prs):
    slide = new_slide(prs, 1, chrome=False)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=DARK, line=None,
         shape_type=MSO_SHAPE.RECTANGLE)
    rect(slide, 0, 0, 0.28, SLIDE_H, fill=YELLOW, line=None,
         shape_type=MSO_SHAPE.RECTANGLE)

    # compass motif (drawn with shapes)
    ring = slide.shapes.add_shape(MSO_SHAPE.DONUT, Inches(9.55), Inches(2.05),
                                  Inches(2.9), Inches(2.9))
    ring.fill.background()
    ring.line.color.rgb = YELLOW
    ring.line.width = Pt(2.0)
    ring.shadow.inherit = False
    style_run(ring.text_frame.paragraphs[0].add_run(), 1, False, DARK)
    try:
        ring.adjustments[0] = 0.03
    except Exception:
        pass

    needle = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(10.70),
                                    Inches(2.62), Inches(0.60), Inches(1.30))
    needle.fill.solid()
    needle.fill.fore_color.rgb = YELLOW
    needle.line.fill.background()
    needle.rotation = 38
    needle.shadow.inherit = False
    style_run(needle.text_frame.paragraphs[0].add_run(), 1, False, DARK)

    needle2 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(10.70),
                                     Inches(3.62), Inches(0.60), Inches(1.30))
    needle2.fill.solid()
    needle2.fill.fore_color.rgb = RGBColor(0x5A, 0x5A, 0x5A)
    needle2.line.fill.background()
    needle2.rotation = 218
    needle2.shadow.inherit = False
    style_run(needle2.text_frame.paragraphs[0].add_run(), 1, False, DARK)

    txbox(slide, 1.10, 1.55, 8.6, 0.4,
          [("2026 금융 AI Challenge  ·  기술설명서", 12, True, YELLOW)])
    txbox(slide, 1.10, 2.05, 8.6, 1.3,
          [("Home_Compass", 46, True, WHITE),
           ("청년 주거 금융 의사결정 서비스", 15, False, GRAY_LT, 8)], line_spacing=1.0)
    rect(slide, 1.14, 3.62, 1.5, 0.05, fill=YELLOW, line=None,
         shape_type=MSO_SHAPE.RECTANGLE)
    txbox(slide, 1.10, 3.88, 8.3, 1.5,
          [("“내가 지금 얼마짜리 집에 살아도 되는가”에", 17, False, WHITE),
           ("3분 만에 숫자로 답하고, 그 결정을 실행할 정책·금융상품까지 연결하는", 17, False, WHITE, 4),
           ("청년 주거 금융 의사결정 에이전트", 17, True, YELLOW, 4)], line_spacing=1.28)

    rect(slide, 1.10, 5.72, 5.1, 0.9, fill=RGBColor(0x38, 0x38, 0x38), line=None,
         radius=0.10, align=PP_ALIGN.LEFT, insets=(0.24, 0.14, 0.20, 0.12),
         anchor=MSO_ANCHOR.MIDDLE,
         items=[("과제 주제", 9.5, True, YELLOW),
                ("「현직자 Pick」 No.1 — 청년 주거 금융 도우미", 12, True, WHITE, 3)])
    rect(slide, 6.42, 5.72, 3.0, 0.9, fill=RGBColor(0x38, 0x38, 0x38), line=None,
         radius=0.10, align=PP_ALIGN.LEFT, insets=(0.24, 0.14, 0.20, 0.12),
         anchor=MSO_ANCHOR.MIDDLE,
         items=[("제출 구분", 9.5, True, YELLOW),
                ("작동형 웹서비스 · 기술설명", 12, True, WHITE, 3)])
    txbox(slide, 1.10, 6.82, 6.0, 0.3,
          [("2026 금융 AI Challenge 제출본", 10.5, False, GRAY_LT)])
    return slide


def slide_02_problem(prs):
    slide = new_slide(
        prs, 2, "청년 주거 의사결정에서 반복되는 4가지 공백",
        "정보는 넘치지만, “나에게 맞는 숫자”를 알려주는 곳은 없습니다.",
        eyebrow="PROBLEM",
    )
    cw, gap = 5.85, 0.33
    x2 = CONTENT_X + cw + gap
    card(slide, CONTENT_X, 1.62, cw, 1.72, "감당 가능 주거비의 기준선이 없다",
         ["소득·자산·부채·가구원수를 넣으면 “월 얼마까지”가 나오는",
          "개인화된 상한선을 제시하는 서비스가 드뭅니다.",
          "결국 “남들 사는 정도”를 기준 삼아 계약합니다."], num="01")
    card(slide, x2, 1.62, cw, 1.72, "전세 vs 월세를 월 납입액으로만 비교한다",
         ["보증금이 묶이며 발생하는 기회비용, 대출이자, 관리비,",
          "보증보험료를 합산한 총비용(TCO) 관점 비교가 부재합니다.",
          "“월세가 싸 보이는 착시”가 그대로 의사결정에 반영됩니다."], num="02")
    card(slide, CONTENT_X, 3.52, cw, 1.72, "정책·상품 정보가 흩어져 있다",
         ["연령·소득·자산·무주택·지역·재직요건이 제도마다 달라",
          "청년이 직접 요건표를 대조해야 적격 여부를 알 수 있습니다.",
          "“되는지 안 되는지” 확인에만 상당한 시간이 듭니다."], num="03")
    card(slide, x2, 3.52, cw, 1.72, "보증금 리스크를 계약 전에 수치로 못 본다",
         ["전세가율, 보증보험 가입 가능 여부 같은 위험 신호를",
          "정량 점수로 확인할 수단이 일반 이용자에게는 부족합니다.",
          "위험은 사후에, 사고가 난 뒤에 학습됩니다."], num="04")

    note_bar(
        slide, 5.52,
        "데이터 정직성 원칙 : 본 장은 검증 불가한 통계 수치를 인용하지 않고, 공개된 제도 요건과 일반적 시장 관행에 근거해 "
        "문제를 정성적으로 기술했습니다.  (참고 성격 자료: 국토교통부 「주거실태조사」, 통계청 「인구주택총조사」)",
        h=0.72,
    )
    txbox(slide, CONTENT_X, 6.36, CONTENT_W, 0.32,
          [("→ 네 가지 공백은 모두 “계산”으로 메울 수 있습니다. 그래서 챗봇이 아니라 판정 엔진을 만들었습니다.",
            12, True, DARK)])
    return slide


def slide_03_gap(prs):
    slide = new_slide(
        prs, 3, "기존 서비스의 한계 — 무엇이 빠져 있는가",
        "네 가지 유형 모두 정보는 주지만, 개인화된 “판정”은 주지 않습니다.",
        eyebrow="AS-IS",
    )
    heads = ["상담형 챗봇", "대출 계산기", "정책 포털·안내서", "부동산 시세 앱"]
    bodies = [
        ["자연어 응답은 매끄럽지만", "수치를 LLM이 직접 생성", "→ 환각·재현성 문제"],
        ["원리금 계산은 정확하지만", "“감당 가능한가” 판단 없음", "→ 기준선 부재"],
        ["제도 목록은 충실하지만", "개인 요건 대조는 사용자 몫", "→ 적격 판정 없음"],
        ["시세·매물은 풍부하지만", "재무 관점 해석이 없음", "→ 총비용 비교 없음"],
    ]
    cw = (CONTENT_W - 3 * 0.28) / 4
    for i in range(4):
        x = CONTENT_X + i * (cw + 0.28)
        rect(slide, x, 1.62, cw, 0.46, fill=DARK, line=None, radius=0.14,
             items=[(heads[i], 12.5, True, WHITE)], insets=(0.10, 0.04, 0.10, 0.04))
        rect(slide, x, 2.08, cw, 1.20, fill=PANEL, line=LINE, radius=0.06,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.20, 0.14, 0.16, 0.08),
             items=[(bodies[i][0], 11, False, DARK_2),
                    (bodies[i][1], 11, False, DARK_2, 4),
                    (bodies[i][2], 11, True, RED, 6)])

    rect(slide, CONTENT_X, 3.50, CONTENT_W, 0.58, fill=WHITE, line=LINE, radius=0.10,
         items=[("공통 공백 :  “정보 제공”은 있으나 “내 상황에 대한 판정”이 없다", 14, True, DARK)])
    arrow(slide, 6.44, 4.16, 0.44, 0.42, "down", fill=YELLOW)

    rect(slide, CONTENT_X, 4.70, CONTENT_W, 1.60, fill=DARK, line=None, radius=0.08,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.34, 0.20, 0.30, 0.16),
         items=[("TO-BE  ·  Home_Compass", 11, True, YELLOW),
                ("① 감당 가능 월주거비 상한을 숫자로 판정하고  ② 전세·월세를 총비용(TCO)과 현재가치(NPV)로 비교하며",
                 12.5, False, WHITE, 8),
                ("③ 제도별 적격 여부를 사유와 함께 판정하고  ④ 보증금 리스크를 0~100점으로 계량합니다.",
                 12.5, False, WHITE, 4),
                ("모든 숫자는 결정론적 엔진이 계산하고, LLM은 그 결과를 설명하는 역할만 맡습니다.",
                 12.5, True, YELLOW, 8)])
    txbox(slide, CONTENT_X, 6.44, CONTENT_W, 0.32,
          [("추천(Recommendation)에서 판정(Assessment)으로 — 이것이 본 과제의 출발점입니다.", 11.5, True, GRAY)])
    return slide


def slide_04_overview(prs):
    slide = new_slide(
        prs, 4, "서비스 개요 — Home_Compass",
        "프로필 입력 3분 → 지불능력·시나리오·정책·리스크를 한 화면에서 판정",
        eyebrow="SOLUTION",
    )
    rect(slide, CONTENT_X, 1.58, CONTENT_W, 1.00, fill=YELLOW_LT, line=None, radius=0.10,
         align=PP_ALIGN.LEFT, insets=(0.34, 0.12, 0.30, 0.12),
         items=[("한 줄 정의", 10, True, YELLOW_DK),
                ("“내가 지금 얼마짜리 집에 살아도 되는가”에 3분 만에 숫자로 답하고,",
                 13.5, True, DARK, 5),
                ("그 결정을 실행할 정책·금융상품까지 연결하는 청년 주거 금융 의사결정 에이전트",
                 13.5, True, DARK, 3)])

    steps = [
        ("STEP 1", "프로필 입력", ["나이 · 연소득 · 월 실수령액", "보유 현금성 자산 · 기존 부채",
                                 "가구원수 · 희망 지역", "무주택 / 신혼 / 재직 요건"]),
        ("STEP 2", "4대 엔진 판정", ["E1 지불능력 상한 산출", "E2 정책 적격성 룰 평가",
                                   "E3 전월세 총비용·NPV 비교", "E4 보증금 리스크 계량"]),
        ("STEP 3", "결과 대시보드", ["감당가능 주거비 상한·권장액", "시나리오 비교 카드 · 적합도 점수",
                                   "적격 정책 리스트 + 판정 사유"]),
        ("STEP 4", "AI 상담", ["자연어 후속 질문", "엔진 결과를 근거로 설명", "숫자는 재계산하지 않음"]),
    ]
    cw = (CONTENT_W - 3 * 0.46) / 4
    for i, (tag, name, lines) in enumerate(steps):
        x = CONTENT_X + i * (cw + 0.46)
        rect(slide, x, 2.76, cw, 1.96, fill=WHITE, line=LINE, radius=0.06,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.22, 0.44, 0.18, 0.12),
             items=[(name, 13.5, True, DARK)] +
                   [(l, 10.5, False, DARK_2, 6 if j == 0 else 3) for j, l in enumerate(lines)])
        badge(slide, x + 0.18, 2.94, 0.78, 0.26, tag, fill=DARK, color=YELLOW, size=8.5)
        if i < 3:
            arrow(slide, x + cw + 0.06, 3.56, 0.34, 0.34, "right", fill=YELLOW)

    outs = [
        # 슈바베지수는 여기 없다. `affordability` 에서 제거되고 시나리오별 실측으로
        # 옮겨졌다 (SPEC 5.2.1 F-1) — 아래 「비교 결과」가 그 자리다.
        ("판정 결과", "감당 가능 월주거비 상한 / 권장액 / 산출 내역 / safe·caution·risk 밴드"),
        ("비교 결과", "시나리오별 5년 총비용, NPV, 월환산비용, 슈바베지수, 적합도 점수"),
        ("연결 결과", "적격 정책·상품 목록과 제도별 판정 사유, 보증금 리스크 점수(0~100)"),
    ]
    ow = (CONTENT_W - 2 * 0.30) / 3
    for i, (t, d) in enumerate(outs):
        x = CONTENT_X + i * (ow + 0.30)
        rect(slide, x, 4.94, ow, 1.24, fill=PANEL, line=LINE, radius=0.06,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.22, 0.16, 0.18, 0.10),
             items=[(t, 11.5, True, YELLOW_DK), (d, 11, False, DARK_2, 5)])
    txbox(slide, CONTENT_X, 6.34, CONTENT_W, 0.34,
          [("모든 출력에는 rationale(판정 근거 문자열) 배열이 함께 반환되어, 화면과 답변에서 근거를 그대로 확인할 수 있습니다.",
            11.5, True, GRAY)])
    return slide


def slide_05_diff(prs):
    slide = new_slide(
        prs, 5, "핵심 차별점 — 우리는 추천하지 않고 판정합니다",
        "“청년 전세대출 추천 챗봇”과 구조적으로 다른 네 가지 설계 결정",
        eyebrow="DIFFERENTIATION",
    )
    diffs = [
        ("D1", "추천이 아니라 판정",
         ["감당 가능 주거비 상한을", "재무적으로 계산해 제시합니다.",
          "“이 상품 어때요”가 아니라", "“당신은 월 OO만원까지”로", "답하는 방식입니다."]),
        ("D2", "총비용·현재가치 비교",
         ["보증금 기회비용까지 반영한", "5년 TCO·NPV로 비교합니다.",
          "월 납입액 착시를 제거하고", "전세와 월세를 하나의 잣대로", "줄 세웁니다."]),
        ("D3", "모든 판정에 근거 동반",
         ["엔진 함수는 예외 없이", "rationale: list[str] 반환.",
          "화면·답변·API 응답 어디서나", "판정 이유를 추적할 수", "있습니다. (XAI)"]),
        ("D4", "환각 차단 + 벤더 비종속",
         ["LLM은 자연어 인터페이스일 뿐,", "숫자는 결정론적 엔진이",
          "100% 계산합니다.", "프로바이더는 교체 가능하고", "키가 없어도 동작합니다."]),
    ]
    cw = (CONTENT_W - 3 * 0.30) / 4
    for i, (tag, title, lines) in enumerate(diffs):
        x = CONTENT_X + i * (cw + 0.30)
        rect(slide, x, 1.62, cw, 2.84, fill=WHITE, line=LINE, radius=0.05,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.24, 0.72, 0.20, 0.14),
             items=[(title, 14, True, DARK)] +
                   [(l, 10.5, False, DARK_2, 8 if j == 0 else 3) for j, l in enumerate(lines)])
        rect(slide, x, 1.62, cw, 0.10, fill=YELLOW, line=None, shape_type=MSO_SHAPE.RECTANGLE)
        badge(slide, x + 0.22, 1.88, 0.56, 0.30, tag, fill=DARK, color=YELLOW, size=10)

    rect(slide, CONTENT_X, 4.70, CONTENT_W, 1.52, fill=NAVY_LT, line=None, radius=0.08,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.32, 0.18, 0.28, 0.12),
         items=[("설계 원칙 한 줄 요약", 10, True, NAVY),
                ("숫자는 코드가 계산하고(결정론), 말은 모델이 하고(자연어), 근거는 항상 함께 반환한다(설명가능성).",
                 14, True, DARK, 7),
                ("이 원칙 때문에 동일 입력에 대해 항상 동일한 판정이 재현되며, 심사·감사·상담 기록으로 그대로 활용할 수 있습니다.",
                 11, False, DARK_2, 6),
                ("또한 LLM 호출이 불가능한 환경에서도 판정 결과가 동일하게 산출되어, 시연·심사 환경에 대한 의존도가 낮습니다.",
                 11, False, DARK_2, 4)])
    txbox(slide, CONTENT_X, 6.36, CONTENT_W, 0.32,
          [("→ 대부분의 유사 과제가 “상담 챗봇”에 머무는 지점에서, 우리는 재현 가능한 판정 엔진을 만들었습니다.",
            11.5, True, GRAY)])
    return slide


def slide_06_architecture(prs):
    slide = new_slide(
        prs, 6, "시스템 아키텍처",
        "4계층 구조 — 클라이언트 / API / 결정론적 판정 엔진 / 데이터, 그리고 교체 가능한 LLM 프로바이더 추상화 계층",
        eyebrow="ARCHITECTURE",
    )
    LX, LW = CONTENT_X, 1.32
    BX = LX + LW + 0.22
    BW = 12.68 - BX

    def layer_label(y, h, num, name):
        rect(slide, LX, y, LW, h, fill=DARK, line=None, radius=0.08,
             align=PP_ALIGN.CENTER,
             items=[(num, 9, True, YELLOW), (name, 11, True, WHITE, 3)],
             insets=(0.06, 0.06, 0.06, 0.06))

    # --- L1 client -------------------------------------------------------
    y1, h1 = 1.50, 0.92
    layer_label(y1, h1, "L1", "클라이언트")
    rect(slide, BX, y1, BW, h1, fill=PANEL, line=LINE, radius=0.05)
    sub_w = (BW - 0.60 - 4 * 0.14) / 3
    subs = [("프로필 입력 폼", "나이·소득·자산·부채·지역"),
            ("결과 대시보드", "인라인 SVG 비교 차트·카드"),
            ("AI 상담 채팅", "자연어 후속 질의")]
    for i, (t, d) in enumerate(subs):
        x = BX + 0.30 + i * (sub_w + 0.14)
        rect(slide, x, y1 + 0.15, sub_w, h1 - 0.30, fill=WHITE, line=LINE, radius=0.06,
             items=[(t, 11, True, DARK), (d, 9, False, GRAY, 3)], insets=(0.08, 0.04, 0.08, 0.04))
    txbox(slide, BX + BW - 3.0, y1 - 0.28, 3.0, 0.26,
          [("Vanilla HTML / CSS / JS  ·  빌드툴·CDN 없음", 9, True, GRAY)], align=PP_ALIGN.RIGHT)

    arrow(slide, 6.44, 2.50, 0.42, 0.30, "down", fill=YELLOW)
    # 정적 `MOCK_RESPONSE` 폴백은 삭제됐다 (SPEC 6.2 D-11 · PR #67 ⑥). 출처 없는 숫자를
    # 조용히 그리는 것이 금지된 것이고, 대체물은 생성물 기반 로컬 엔진 + 화면 명시다.
    txbox(slide, 7.00, 2.50, 5.6, 0.30,
          [("HTTP / JSON  ·  fetch()  ·  백엔드 미기동 시 생성물 기반 로컬 판정", 9.5, False, GRAY)])

    # --- L2 API ----------------------------------------------------------
    y2, h2 = 2.88, 0.84
    layer_label(y2, h2, "L2", "API 계층")
    rect(slide, BX, y2, BW, h2, fill=NAVY_LT, line=LINE, radius=0.05)
    # 파일 이름은 `main.py` 다. 그리고 CORS 는 **없다** — `allow_origins=["*"]` 는
    # 쿠키 세션과 양립하지 않아 폐기됐고, web·admin·api 를 같은 오리진에서 서빙한다
    # (SPEC 6.2 · main.py 의 StaticFiles 마운트).
    txbox(slide, BX + 0.28, y2 + 0.10, 3.1, 0.30, [("FastAPI  ·  main.py", 11.5, True, NAVY)])
    txbox(slide, BX + 0.28, y2 + 0.44, 3.4, 0.30,
          [("동일 오리진 서빙 · 쿠키 세션", 9, False, GRAY)])
    eps = ["GET /api/health", "GET /api/regions", "POST /api/analyze", "POST /api/chat"]
    ep_w = 1.49
    for i, ep in enumerate(eps):
        rect(slide, BX + 3.70 + i * (ep_w + 0.13), y2 + 0.22, ep_w, 0.40,
             fill=WHITE, line=LINE, radius=0.14,
             items=[(ep, 8.5, True, DARK_2)], insets=(0.03, 0.02, 0.03, 0.02))
    # 이 네 칸은 **시민 판정 경로**만이다. 개수를 적지 않으면 넷이 전부로 읽힌다.
    # ★ 수를 손으로 적지 않는다 — PR #76 이 「15종」을 적어 둔 뒤 6-A 신고 둘과 7단계
    #   감사 하나가 늘어 **18종이 됐고 아무도 못 잡았다.** `contracts/openapi.json` 을
    #   세어서 쓴다 (`_endpoint_counts()`). 계약이 곧 정본이다.
    _cite, _cadmin, _ctotal = _endpoint_counts()
    txbox(slide, BX + 3.70, y2 + 0.63, 6.35, 0.20,
          [("시민 판정 경로 4종  ·  계약 전체는 %d종 (인증 3 · 규칙관리 %d · meta 1 포함)"
            % (_ctotal, _cadmin), 8, False, GRAY)])

    arrow(slide, 4.30, 3.80, 0.42, 0.30, "down", fill=YELLOW)
    arrow(slide, 8.60, 3.80, 0.42, 0.30, "down", fill=YELLOW)

    # --- L3 engines ------------------------------------------------------
    y3, h3 = 4.18, 1.30
    layer_label(y3, h3, "L3", "판정 엔진")
    eng_area_w = BW - 3.55
    rect(slide, BX, y3, eng_area_w, h3, fill=YELLOW_LT, line=LINE, radius=0.05)
    engines = [("E1", "주거지불능력", "Affordability"), ("E2", "정책 적격성", "Eligibility"),
               ("E3", "총비용·NPV", "TCO / NPV"), ("E4", "보증금 리스크", "Risk Scanner")]
    ew = (eng_area_w - 0.44 - 3 * 0.14) / 4
    for i, (tag, ko, en) in enumerate(engines):
        x = BX + 0.22 + i * (ew + 0.14)
        rect(slide, x, y3 + 0.18, ew, h3 - 0.36, fill=WHITE, line=YELLOW, line_w=1.4,
             radius=0.07,
             items=[(tag, 12, True, YELLOW_DK), (ko, 11.5, True, DARK, 3),
                    (en, 8.5, False, GRAY, 2)], insets=(0.05, 0.05, 0.05, 0.05))
    txbox(slide, BX + 0.22, y3 + h3 + 0.02, 6.4, 0.26,
          [("순수 함수 · 동일 입력 → 동일 출력 · 모든 반환값에 rationale 포함", 9, False, GRAY)])

    # LLM provider-abstraction layer at the right, coupled to engines.
    # The 3-tier fallback ladder is drawn as stacked bars + a downward arrow.
    ax = BX + eng_area_w + 0.28
    aw = 12.68 - ax
    rect(slide, ax, y3, aw, h3, fill=DARK, line=None, radius=0.06)
    txbox(slide, ax + 0.16, y3 + 0.09, aw - 0.30, 0.24,
          [("A1  LLM 프로바이더 추상화", 10.5, True, YELLOW)])
    tiers = [("① OpenAI function calling", YELLOW, DARK),
             ("② Anthropic tool use", RGBColor(0x6A, 0x6A, 0x6A), WHITE),
             ("③ 룰 기반 결정론적 폴백", GREEN, WHITE)]
    bar_x, bar_w = ax + 0.44, aw - 0.60
    for i, (label, fill_c, txt_c) in enumerate(tiers):
        rect(slide, bar_x, y3 + 0.38 + i * 0.29, bar_w, 0.25, fill=fill_c, line=None,
             radius=0.30, items=[(label, 8, True, txt_c)],
             insets=(0.06, 0.01, 0.04, 0.01), align=PP_ALIGN.LEFT)
    conn(slide, ax + 0.28, y3 + 0.40, ax + 0.28, y3 + 1.14, color=YELLOW, width=1.4)
    txbox(slide, ax, y3 + h3 + 0.02, aw + 0.10, 0.26,
          [("키 미설정·호출 실패 시 자동 강등", 8.5, False, GRAY)])
    conn(slide, ax, y3 + 0.52, ax - 0.28, y3 + 0.52, color=DARK, width=1.6)
    conn(slide, ax - 0.28, y3 + 0.82, ax, y3 + 0.82, color=DARK, width=1.6)

    arrow(slide, 6.44, 5.58, 0.42, 0.28, "down", fill=YELLOW)

    # --- L4 data ---------------------------------------------------------
    y4, h4 = 5.94, 0.78
    layer_label(y4, h4, "L4", "데이터")
    rect(slide, BX, y4, BW, h4, fill=PANEL, line=LINE, radius=0.05)
    # 런타임 판정은 **저장소**에서 읽는다 (`store.regions.list()` · `RuleVersion.payload`).
    # `data/*.json` 은 그 저장소를 채우는 **시드 입력**이지 판정이 읽는 곳이 아니다
    # (SPEC 5.1.1 fail-closed · scripts/seed_store.py).
    data_items = [("저장소 (SQLite)", "판정이 읽는 곳 · 지역 10 · 규칙 8"),
                  ("data/*.json", "저장소를 채우는 시드 입력"),
                  ("source / 고지", "정책 항목마다 출처·고지 필드 필수")]
    dw = (BW - 0.60 - 2 * 0.14) / 3
    for i, (t, d) in enumerate(data_items):
        x = BX + 0.30 + i * (dw + 0.14)
        rect(slide, x, y4 + 0.12, dw, h4 - 0.24, fill=WHITE, line=LINE, radius=0.06,
             items=[(t, 10.5, True, DARK), (d, 9, False, GRAY, 2)],
             insets=(0.06, 0.03, 0.06, 0.03))
    return slide


def slide_07_pipeline(prs):
    slide = new_slide(
        prs, 7, "4대 엔진 처리 흐름",
        "지불능력 상한을 먼저 확정한 뒤, 그 상한을 제약조건으로 정책·시나리오·리스크를 평가합니다.",
        eyebrow="ENGINE PIPELINE",
    )
    ws = [2.55, 2.65, 2.75, 2.58]
    xs = [CONTENT_X]
    for i in range(3):
        xs.append(xs[-1] + ws[i] + 0.50)

    ROW_Y, ROW_H = 2.56, 1.62
    ROW_MID = ROW_Y + ROW_H / 2

    # stage 1 : input
    rect(slide, xs[0], ROW_Y, ws[0], ROW_H, fill=DARK, line=None, radius=0.06,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.20, 0.16, 0.16, 0.10),
         items=[("INPUT", 9, True, YELLOW),
                ("사용자 프로필", 13, True, WHITE, 3),
                ("나이 · 연소득 · 월 실수령액", 9.5, False, GRAY_LT, 7),
                ("자산 · 기존부채 · 가구원수", 9.5, False, GRAY_LT, 3),
                ("희망지역 · 무주택 · 재직요건", 9.5, False, GRAY_LT, 3)])
    arrow(slide, xs[0] + ws[0] + 0.06, ROW_MID - 0.15, 0.38, 0.30, "right", fill=YELLOW)

    # stage 2 : E1
    rect(slide, xs[1], ROW_Y, ws[1], ROW_H, fill=WHITE, line=YELLOW, line_w=1.6, radius=0.06,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.20, 0.16, 0.16, 0.10),
         items=[("E1  주거지불능력", 12.5, True, YELLOW_DK),
                ("감당 가능 월주거비 상한", 10.5, True, DARK, 6),
                ("· 권장 주거비 · 산출 내역", 9.5, False, DARK_2, 5),
                ("· safe / caution / risk 밴드", 9.5, False, DARK_2, 3),
                ("이후 모든 판정의 제약조건", 9, True, GRAY, 6)])

    # branch from E1 to E2 / E3
    bx = xs[1] + ws[1]
    conn(slide, bx + 0.06, ROW_MID, xs[2] - 0.04, 2.26, color=YELLOW, width=2.0)
    conn(slide, bx + 0.06, ROW_MID, xs[2] - 0.04, 4.52, color=YELLOW, width=2.0)

    # stage 3 : E2 / E3 (parallel)
    rect(slide, xs[2], 1.66, ws[2], 1.20, fill=PANEL, line=LINE, radius=0.06,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.20, 0.14, 0.16, 0.08),
         items=[("E2  정책·상품 적격성", 12, True, DARK),
                ("제도별 eligible / conditional /", 9.5, False, DARK_2, 6),
                ("ineligible + 판정 사유 배열", 9.5, False, DARK_2, 3)])
    rect(slide, xs[2], 3.92, ws[2], 1.20, fill=PANEL, line=LINE, radius=0.06,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.20, 0.14, 0.16, 0.08),
         items=[("E3  전월세 총비용 비교", 12, True, DARK),
                ("시나리오별 5년 TCO · NPV ·", 9.5, False, DARK_2, 6),
                ("월환산비용 · 적합도 점수", 9.5, False, DARK_2, 3)])
    txbox(slide, xs[2], 3.22, ws[2], 0.28,
          [("병렬 평가", 9.5, True, GRAY)], align=PP_ALIGN.CENTER)

    # merge into E4
    ex = xs[2] + ws[2]
    conn(slide, ex + 0.04, 2.26, xs[3] - 0.06, ROW_MID - 0.18, color=YELLOW, width=2.0)
    conn(slide, ex + 0.04, 4.52, xs[3] - 0.06, ROW_MID + 0.18, color=YELLOW, width=2.0)

    # stage 4 : E4
    rect(slide, xs[3], ROW_Y, ws[3], ROW_H, fill=WHITE, line=YELLOW, line_w=1.6, radius=0.06,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.20, 0.16, 0.16, 0.10),
         items=[("E4  보증금 리스크", 12.5, True, YELLOW_DK),
                ("위험 점수 0~100", 10.5, True, DARK, 6),
                ("· low / medium / high 밴드", 9.5, False, DARK_2, 5),
                ("· 위험요인별 영향도 목록", 9.5, False, DARK_2, 3),
                # engines/__init__.py 는 `scenarios[0]`(적합도 1위)의 보증금으로 돈다.
                # 그 1위가 월세면 월세 보증금으로 돈다 — 전세 전용이 아니다.
                ("적합도 1위 안의 보증금에 적용", 9, True, GRAY, 6)])

    conn(slide, xs[3] + ws[3] / 2, ROW_Y + ROW_H, xs[3] + ws[3] / 2, 5.10,
         color=YELLOW, width=2.0)

    # output band
    rect(slide, CONTENT_X, 5.16, 8.60, 1.22, fill=YELLOW_LT, line=None, radius=0.07,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.28, 0.16, 0.24, 0.10),
         items=[("OUTPUT  ·  종합 판정 리포트  (POST /api/analyze 단일 응답)", 11.5, True, YELLOW_DK),
                ("affordability · scenarios[] · policies[] · risk · summary · meta"
                 " · internal (상담원 로그인 시)",
                 11, True, DARK, 7),
                ("모든 블록에 rationale / reasons / factors 배열이 함께 담겨 판정 근거를 그대로 추적할 수 있습니다.",
                 10, False, DARK_2, 5)])
    arrow(slide, 9.00, 5.62, 0.34, 0.30, "right", fill=YELLOW)
    rect(slide, 9.46, 5.16, 3.22, 1.22, fill=DARK, line=None, radius=0.07,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.22, 0.16, 0.18, 0.10),
         items=[("A1  LLM 프로바이더 추상화", 11, True, YELLOW),
                ("판정 결과를 근거로 자연어 설명", 10, False, WHITE, 7),
                ("숫자 재계산 없음 (POST /api/chat)", 9.5, False, GRAY_LT, 5)])
    return slide


def _engine_header(slide, tag, ko, en, one_liner):
    rect(slide, CONTENT_X, 1.56, CONTENT_W, 0.56, fill=DARK, line=None, radius=0.08,
         align=PP_ALIGN.LEFT, insets=(0.26, 0.06, 0.22, 0.06),
         items=[("%s  ·  %s  (%s)   —   %s" % (tag, ko, en, one_liner), 12, True, WHITE)])


def _io_row(slide, y, h, in_lines, proc_lines, out_lines):
    w_in, w_proc, w_out = 3.45, 3.75, 3.43
    x_in = CONTENT_X
    x_proc = x_in + w_in + 0.72
    x_out = x_proc + w_proc + 0.68

    def block(x, w, tag, color, lines):
        rect(slide, x, y, w, h, fill=WHITE, line=LINE, radius=0.06,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.24, 0.42, 0.20, 0.12),
             items=[(l, 10.5, False, DARK_2, 0 if i == 0 else 4) for i, l in enumerate(lines)])
        badge(slide, x + 0.20, y + 0.14, 0.86, 0.26, tag, fill=color,
              color=WHITE if color != YELLOW else DARK, size=8.5)

    block(x_in, w_in, "INPUT", NAVY, in_lines)
    block(x_proc, w_proc, "PROCESS", YELLOW, proc_lines)
    block(x_out, w_out, "OUTPUT", GREEN, out_lines)
    arrow(slide, x_in + w_in + 0.16, y + h / 2 - 0.16, 0.40, 0.32, "right", fill=YELLOW)
    arrow(slide, x_proc + w_proc + 0.14, y + h / 2 - 0.16, 0.40, 0.32, "right", fill=YELLOW)


def _rationale_box(slide, y, lines, h=1.06):
    rect(slide, CONTENT_X, y, CONTENT_W, h, fill=YELLOW_LT, line=None, radius=0.08,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.30, 0.12, 0.26, 0.10),
         items=[("rationale 예시  (아래 문자열은 엔진이 실제로 반환하는 판정 근거의 형식 예시입니다)",
                 9.5, True, YELLOW_DK)] +
               [("“%s”" % l, 11, False, DARK, 6 if i == 0 else 3)
                for i, l in enumerate(lines)])


def slide_08_e1(prs):
    slide = new_slide(prs, 8, "E1 · 주거지불능력 엔진 (Affordability)",
                      "“월 얼마까지 감당 가능한가”를 가처분소득 기반으로 판정합니다.",
                      eyebrow="ENGINE 1 / 4")
    _engine_header(slide, "E1", "주거지불능력 엔진", "Affordability",
                   "감당 가능 월주거비 상한과 권장액을 산출하는 기준선 엔진")
    # PROCESS 는 엔진의 실제 순서다 (engines/affordability.py). 상한은 **두 개를 계산해
    # 작은 쪽**을 쓴다 — 그 규칙이 이 엔진의 요지이고, 예전 문구("가처분소득 = 실수령액 −
    # 기존부채")는 잔여여력 경로에서만 하는 차감을 가처분소득의 정의로 적어 틀렸다.
    # OUTPUT 에 `schwabeIndexPct` 는 없다 (SPEC 5.2.1 F-1).
    _io_row(slide, 2.20, 1.48,
            ["월 실수령액 (monthlyNetIncomeKRW)", "연소득 (annualIncomeKRW)",
             "기존 부채 월 상환액", "가구원수 (householdSize)"],
            ["① 비율상한 = 가처분소득 × 상한비율", "② 잔여여력 = 소득 − 생활비 − 부채 − 버퍼",
             "③ 상한 = 두 값 중 작은 쪽", "④ 권장액 = 상한 × 안전마진"],
            ["maxMonthlyHousingCostKRW", "recommendedMonthlyHousingCostKRW",
             "breakdown (항목별 분해)", "band · rationale[]"])

    # band gauge drawn with shapes
    # 밴드는 슈바베지수 구간이 아니다. 상한/소득(cap_ratio)과 부채/소득(debt_ratio)
    # **두 축**을 함께 본다 (engines/affordability.py).
    txbox(slide, CONTENT_X, 3.80, 7.30, 0.30,
          [("판정 밴드 — 상한÷소득(cap_ratio)과 부채÷소득(debt_ratio) 두 축으로 판정", 11, True, DARK)])
    gx, gw, gy, gh = CONTENT_X, 7.30, 4.16, 0.46
    segs = [("safe", GREEN, GREEN_LT, 0.34, "여유 있음"),
            ("caution", AMBER, YELLOW_LT, 0.33, "점검 필요"),
            ("risk", RED, RED_LT, 0.33, "과부담 위험")]
    cx = gx
    for name, fg, bg, frac, cap in segs:
        w = gw * frac
        rect(slide, cx, gy, w, gh, fill=bg, line=fg, line_w=1.2, radius=0.20,
             items=[(name, 11, True, fg)], insets=(0.04, 0.02, 0.04, 0.02))
        txbox(slide, cx, gy + gh + 0.05, w, 0.26, [(cap, 9, False, GRAY)],
              align=PP_ALIGN.CENTER)
        cx += w + 0.04
    rect(slide, 8.20, 4.10, 4.48, 0.98, fill=PANEL, line=LINE, radius=0.06,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.22, 0.14, 0.18, 0.08),
         items=[("breakdown 반환 항목", 10.5, True, DARK),
                ("netIncome · livingCost · existingDebt · buffer", 10, False, DARK_2, 6),
                ("→ 상한이 왜 그 값인지 항목별로 분해합니다.", 9.5, False, GRAY, 4)])

    # 아래 두 줄은 지어낸 예시가 아니라 POST /api/analyze 가 실제로 반환한 문자열이다.
    # 예전 두 번째 줄("상한 대비 실제 주거비 비중이 안전 구간")은 존재하지 않는 판정
    # 경로를 설명하고 있었다 — 밴드는 cap_ratio·debt_ratio 로 갈린다.
    _rationale_box(slide, 5.26,
                   ["소득 대비 비율보다 실제 잔여 여력이 더 빠듯해, 잔여 여력 기준으로 상한을 86만원으로 확정했습니다.",
                    "소득 대비 주거비 여력이 안정적인 구간(safe)입니다."])
    txbox(slide, CONTENT_X, 6.44, CONTENT_W, 0.30,
          [("경계값(소득 0원, 부채 과다, 고소득) 케이스는 pytest 단위 테스트로 검증합니다.", 10, False, GRAY)])
    return slide


def slide_09_e2(prs):
    slide = new_slide(prs, 9, "E2 · 정책·상품 적격성 룰엔진 (Eligibility)",
                      "제도별 요건을 순차 평가해 적격 여부와 그 사유를 함께 반환합니다.",
                      eyebrow="ENGINE 2 / 4")
    _engine_header(slide, "E2", "정책·상품 적격성 룰엔진", "Eligibility",
                   "LLM 추론이 아닌 명시적 규칙으로 적격 여부를 판정하는 룰엔진")
    _io_row(slide, 2.20, 1.48,
            ["나이 · 연소득 · 보유자산", "무주택 여부 (isHomeless)",
             "희망 지역코드 · 신혼 여부", "중소기업 재직 여부"],
            # ★ 판정이 읽는 것은 `policies.json` 이 **아니다.** 저장소의 승인된
            #   `RuleVersion` 을 읽는다 (`main.read_active_rule_versions`). `data/*.json`
            #   은 그 저장소를 채우는 시드 입력이다. 슬라이드 06 은 PR #76 이 이미
            #   그렇게 고쳤는데 이 칸이 안 따라와 **덱이 자기모순**이었다.
            ["저장소의 승인된 규칙(RuleVersion)을 순차 평가", "요건 전부 충족 → eligible",
             "일부 미확인·조건부 → conditional", "명확한 미충족 → ineligible"],
            ["status (3값) · reasons[]", "maxAmountKRW · rateRangePct",
             "source (출처 기관명)", "disclaimer (고지 문구)"])

    # rule chips -> verdict boxes
    txbox(slide, CONTENT_X, 3.78, 6.0, 0.30, [("평가 규칙 항목", 11, True, DARK)])
    # policies.json 의 `criteria` 축 전수 — 7종이다. 예전 목록은 신혼(requireNewlywed)
    # 을 빠뜨렸고, 그 축이 실제로 한 정책을 ineligible 로 떨군다.
    chips = ["연령 요건", "소득 요건", "자산 요건", "무주택 요건", "지역 요건", "재직 요건", "혼인 요건"]
    cw = (6.85 - 6 * 0.10) / 7
    for i, ch in enumerate(chips):
        rect(slide, CONTENT_X + i * (cw + 0.10), 4.12, cw, 0.42,
             fill=NAVY_LT, line=NAVY, line_w=1.0, radius=0.24,
             items=[(ch, 9.5, True, NAVY)], insets=(0.03, 0.02, 0.03, 0.02))
    arrow(slide, 7.58, 4.18, 0.36, 0.30, "right", fill=YELLOW)

    verdicts = [("eligible", "요건 전부 충족", GREEN, GREEN_LT),
                ("conditional", "추가 확인 필요", AMBER, YELLOW_LT),
                ("ineligible", "요건 미충족", RED, RED_LT)]
    vw = (12.68 - 8.06 - 2 * 0.10) / 3
    for i, (name, desc, fg, bg) in enumerate(verdicts):
        rect(slide, 8.06 + i * (vw + 0.10), 4.04, vw, 0.58, fill=bg, line=fg, line_w=1.2,
             radius=0.10, items=[(name, 10.5, True, fg), (desc, 8.5, False, DARK_2, 2)],
             insets=(0.04, 0.03, 0.04, 0.03))

    rect(slide, CONTENT_X, 4.74, CONTENT_W, 0.50, fill=PANEL, line=LINE, radius=0.08,
         align=PP_ALIGN.LEFT, insets=(0.26, 0.06, 0.22, 0.06),
         items=[("판정은 항상 사유와 1:1로 대응합니다 — reasons[] 길이가 0인 판정은 반환하지 않습니다.",
                 11, True, DARK)])

    # POST /api/analyze 가 실제로 반환한 reasons[] 에서 가져왔다 (eligible 1줄 + conditional 1줄).
    _rationale_box(slide, 5.34,
                   ["만 19~34세 요건 충족 (28세) · 연소득 5,000만원 이하 충족 (4,200만원)",
                    "주택 유형·전세가율·선순위 채권 규모에 따라 보증 가입이 제한될 수 있습니다. (conditional)"])
    txbox(slide, CONTENT_X, 6.50, CONTENT_W, 0.30,
          [("정책 요건은 공개된 일반 기준 수준으로만 기술하며, 구체적 금리·한도 수치는 예시값으로 명시합니다. (18장 참조)",
            10, False, GRAY)])
    return slide


def slide_10_e3(prs):
    slide = new_slide(prs, 10, "E3 · 전월세 총비용 비교 엔진 (TCO / NPV)",
                      "월 납입액이 아니라 5년 총비용과 현재가치로 시나리오를 동일 기준 비교합니다.",
                      eyebrow="ENGINE 3 / 4")
    _engine_header(slide, "E3", "전월세 총비용 비교 엔진", "TCO / NPV",
                   "보증금 기회비용까지 반영해 전세·월세를 하나의 잣대로 세우는 비교 엔진")
    _io_row(slide, 2.20, 1.48,
            ["보증금 (depositKRW)", "월세 · 관리비", "대출금액 · 대출금리", "보유 자산 · 비교 기간(5년)"],
            ["① 비용 구성요소를 월 단위로 전개", "② 5년 누적 → tco5yKRW",
             "③ 할인율 적용 → npv5yKRW", "④ E1 상한과 대조 → verdict·fitScore"],
            ["monthlyEquivalentCostKRW", "tco5yKRW · npv5yKRW",
             "components (항목별 분해)", "fitScore · verdict · rationale[]"])

    # cost component composition drawn as a stacked bar
    txbox(slide, CONTENT_X, 3.76, 7.0, 0.30,
          [("총비용 구성요소 — 월세만 보면 보이지 않는 항목들", 11, True, DARK)])
    comps = [("대출이자", 0.20, NAVY), ("월세", 0.22, YELLOW_DK), ("관리비", 0.16, GRAY),
             ("보증금 기회비용", 0.26, RED), ("보증보험료", 0.16, GREEN)]
    bx, bw, by, bh = CONTENT_X, 7.30, 4.12, 0.50
    cx = bx
    for name, frac, col in comps:
        w = bw * frac
        rect(slide, cx, by, w, bh, fill=col, line=None, radius=0.10,
             items=[(name, 8.5, True, WHITE)], insets=(0.02, 0.02, 0.02, 0.02))
        cx += w + 0.03
    txbox(slide, bx, by + bh + 0.06, bw, 0.30,
          [("※ 막대 길이는 구성요소를 설명하기 위한 도식이며, 실제 비중은 입력값에 따라 계산됩니다.",
            8.5, False, GRAY)])

    rect(slide, 8.20, 3.76, 4.48, 1.46, fill=PANEL, line=LINE, radius=0.06,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.22, 0.16, 0.18, 0.08),
         items=[("핵심 아이디어 — 기회비용", 11, True, DARK),
                ("전세 보증금은 “사라지는 돈”은 아니지만,", 10, False, DARK_2, 5),
                ("그 기간 동안 다른 곳에 쓰일 수 없습니다.", 10, False, DARK_2, 3),
                ("이 묶임 비용을 총비용에 포함해야", 10, True, DARK, 4),
                ("전세와 월세를 공정하게 비교할 수 있습니다.", 10, True, DARK, 3)])

    # POST /api/analyze 의 scenarios[].rationale 에서 그대로 가져왔다. 예전 두 번째 줄이
    # 든 예(권장액 이내 → affordable)는 이 프로필에서 나오지 않는 분기였다.
    _rationale_box(slide, 5.30,
                   ["보증금에 묶이는 자기자본 8,000만원의 기회비용을 연 3.3%로 계산해 1,304만원을 총비용에 포함했습니다.",
                    "월 환산비용이 권장액 73만원은 넘지만 상한 86만원 이내여서 무리하면 가능한 구간입니다. (stretch)"])
    txbox(slide, CONTENT_X, 6.46, CONTENT_W, 0.30,
          [("비교 기간과 할인율은 화면에서 확인 가능한 가정값으로 노출하며, 가정을 바꾸면 결과도 함께 갱신됩니다.",
            10, False, GRAY)])
    return slide


def slide_11_e4(prs):
    slide = new_slide(prs, 11, "E4 · 전세보증금 리스크 스캐너 (Risk)",
                      "계약 전에 확인 가능한 위험 신호를 0~100점으로 계량해 보여줍니다.",
                      eyebrow="ENGINE 4 / 4")
    _engine_header(slide, "E4", "전세보증금 리스크 스캐너", "Risk Scanner",
                   "위험을 사후 학습이 아니라 계약 전 숫자로 확인하게 하는 스캐너")
    # scan_deposit_risk 의 실제 인자다. 「보유자산」은 이 엔진에 들어오지 않는다 —
    # 비중을 재는 상대는 보유자산이 아니라 **보증금**이고(대출/보증금), 지역은
    # regions.json 의 `marketRisk` 로 들어온다.
    _io_row(slide, 2.20, 1.48,
            ["보증금 (depositKRW)", "전세가율 (jeonseRatioPct)",
             "대출금액 · 보증보험 가입 가능 여부", "지역 시장 위험도 (marketRisk)"],
            ["① 위험요인별 값을 정규화", "② 요인별 가중치 적용 후 합산",
             "③ 0~100 점수로 스케일링", "④ low / medium / high 밴드 판정"],
            ["score (0~100)", "band (low / medium / high)",
             "factors[] : name · valuePct", "factors[].impact · note"])

    txbox(slide, CONTENT_X, 3.76, 7.0, 0.30, [("위험 점수 밴드", 11, True, DARK)])
    gx, gw, gy, gh = CONTENT_X, 7.30, 4.10, 0.44
    # 경계는 `risk.band_low_max` / `risk.band_medium_max` 다 (현행 34 / 64).
    # 예전 값 39 / 69 는 어느 시점의 것도 아니게 됐다 — 상수 레지스트리가 정본이므로
    # 아래 캡션이 그 사실을 함께 적는다. 슬라이드가 두 번째 정본이 되면 또 갈라진다.
    bands = [("low  0–34", GREEN, GREEN_LT, 0.40), ("medium  35–64", AMBER, YELLOW_LT, 0.30),
             ("high  65–100", RED, RED_LT, 0.30)]
    cx = gx
    for name, fg, bg, frac in bands:
        w = gw * frac
        rect(slide, cx, gy, w, gh, fill=bg, line=fg, line_w=1.2, radius=0.20,
             items=[(name, 10, True, fg)], insets=(0.03, 0.02, 0.03, 0.02))
        cx += w + 0.04
    txbox(slide, gx, gy + gh + 0.05, gw, 0.28,
          [("※ 경계값의 정본은 상수 레지스트리(model_constants)이며, 실제 심사 기준과는 무관합니다.",
            8.5, False, GRAY)])

    # engines/risk.py 가 실제로 담는 factors[].name 전수 — 5개다.
    factors = ["전세가율", "보증보험 가입 가능성", "보증금 내 대출 비중",
               "보증금 규모 (상한 대비)", "지역 임대차 시장 여건"]
    fw = (7.30 - 4 * 0.10) / 5
    for i, f in enumerate(factors):
        rect(slide, gx + i * (fw + 0.10), 4.86, fw, 0.42, fill=WHITE, line=LINE, radius=0.22,
             items=[(f, 9.5, True, DARK_2)], insets=(0.03, 0.02, 0.03, 0.02))

    rect(slide, 8.20, 3.76, 4.48, 1.46, fill=PANEL, line=LINE, radius=0.06,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.22, 0.16, 0.18, 0.08),
         items=[("설계 의도", 11, True, DARK),
                ("점수는 “안전 보증”이 아니라", 10, False, DARK_2, 5),
                ("“확인해야 할 항목의 우선순위”입니다.", 10, False, DARK_2, 3),
                ("각 요인에 impact(low/medium/high)와", 10, False, DARK_2, 4),
                ("note를 붙여 점검 항목을 안내합니다.", 10, False, DARK_2, 3)])

    # POST /api/analyze 의 risk.rationale 에서 그대로 가져왔다.
    _rationale_box(slide, 5.38,
                   ["서울 마포구 기준 보증금 2억 8,000만원에 대한 위험 점수는 100점 만점에 32점(low)입니다.",
                    "가장 큰 위험 요인은 '보증금 내 대출 비중'입니다: 보증금의 70% 이상이 대출로, 보증금 사고 시 부채만 남을 위험이 있습니다."])
    txbox(slide, CONTENT_X, 6.54, CONTENT_W, 0.30,
          [("본 엔진은 법적·계약적 안전을 보증하지 않으며, 계약 전 점검 항목을 우선순위화하는 참고 지표입니다.",
            10, True, GRAY)])
    return slide


def slide_12_llm(prs):
    slide = new_slide(prs, 12, "A1 · LLM 프로바이더 추상화와 환각 차단 설계",
                      "LLM은 말을 하고, 숫자는 엔진이 만듭니다 — 그리고 그 LLM은 언제든 교체할 수 있습니다.",
                      eyebrow="LLM ABSTRACTION & GUARDRAIL")

    # --- request flow -----------------------------------------------------
    steps = [("사용자 질문", "“전세랑 월세 중\n뭐가 나아?”"),
             ("추상화 계층", "프로바이더 선택 ·\n도구 스키마 변환"),
             ("도구 호출", "E1~E4 결정론적\n엔진 함수 호출"),
             ("엔진 반환", "숫자 + rationale\n(재계산 불가)"),
             ("답변 생성", "반환값만 인용해\n문장 구성")]
    sw = (CONTENT_W - 4 * 0.42) / 5
    for i, (t, d) in enumerate(steps):
        x = CONTENT_X + i * (sw + 0.42)
        fill = DARK if i in (0, 4) else WHITE
        line_c = None if i in (0, 4) else YELLOW
        tcol = WHITE if i in (0, 4) else DARK
        dcol = GRAY_LT if i in (0, 4) else DARK_2
        lines = d.split("\n")
        rect(slide, x, 1.54, sw, 1.06, fill=fill, line=line_c, line_w=1.5, radius=0.07,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.18, 0.12, 0.14, 0.06),
             items=[(t, 11.5, True, tcol)] +
                   [(l, 9.5, False, dcol, 5 if j == 0 else 2) for j, l in enumerate(lines)])
        if i < 4:
            arrow(slide, x + sw + 0.04, 1.92, 0.34, 0.30, "right", fill=YELLOW)

    # --- 3-tier provider fallback ladder ---------------------------------
    txbox(slide, CONTENT_X, 2.70, 9.5, 0.28,
          [("① LLM 프로바이더 추상화 — 3단 폴백 구조 (동일한 도구 정의를 프로바이더별 규격으로 변환)",
            11.5, True, DARK)])
    tiers = [
        ("1순위", "OpenAI function calling", YELLOW, DARK, YELLOW_LT,
         ["tools 스키마로 E1~E4를 노출", "OPENAI_API_KEY 설정 시 사용"]),
        ("2순위", "Anthropic tool use", GRAY, WHITE, PANEL,
         ["동일 도구 정의를 자동 변환", "ANTHROPIC_API_KEY 설정 시 사용"]),
        ("3순위", "룰 기반 결정론적 폴백", GREEN, WHITE, GREEN_LT,
         ["키 없음·호출 실패 시 자동 전환", "엔진 결과를 템플릿 문장으로 반환"]),
    ]
    tw = (CONTENT_W - 2 * 0.62) / 3
    for i, (rank, name, accent, rank_txt, bg, lines) in enumerate(tiers):
        x = CONTENT_X + i * (tw + 0.62)
        rect(slide, x, 3.00, tw, 1.26, fill=bg, line=accent, line_w=1.4, radius=0.06,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.22, 0.42, 0.18, 0.08),
             items=[(name, 12, True, DARK)] +
                   [(l, 9.5, False, DARK_2, 6 if j == 0 else 3) for j, l in enumerate(lines)])
        badge(slide, x + 0.20, 3.14, 0.62, 0.26, rank, fill=accent, color=rank_txt, size=8.5)
        if i < 2:
            txbox(slide, x + tw - 0.06, 3.20, 0.74, 0.26,
                  [("실패 시", 8, True, GRAY)], align=PP_ALIGN.CENTER)
            arrow(slide, x + tw + 0.08, 3.50, 0.46, 0.30, "right", fill=YELLOW)

    rect(slide, CONTENT_X, 4.36, CONTENT_W, 0.68, fill=YELLOW_LT, line=None, radius=0.10,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.24, 0.10, 0.22, 0.06),
         items=[("강점 : 특정 벤더에 종속되지 않으며, 금융권 도입 시 내부망 전용 모델로 프로바이더만 교체하면 됩니다.",
                 10.5, False, DARK_2),
                ("키가 하나도 없어도 3순위 폴백으로 핵심 판정 기능은 100% 동작합니다. (graceful degradation)",
                 10.5, False, DARK_2, 4)])

    # --- hallucination guardrails ----------------------------------------
    txbox(slide, CONTENT_X, 5.12, 9.5, 0.28,
          [("② 환각 차단 3중 장치 — 프로바이더가 무엇이든 숫자는 바뀌지 않습니다", 11.5, True, DARK)])
    guards = [
        ("G1  숫자 생성 권한 분리",
         ["금액·비율·점수는 엔진 반환값만 사용하고,", "모델이 산출한 수치는 답변에 싣지 않습니다."]),
        ("G2  근거 문자열 강제 동반",
         ["모든 엔진이 rationale을 반환하고,", "답변은 그 문자열을 근거로 구성됩니다."]),
        ("G3  프로바이더 무관 재현성",
         ["1·2·3순위 어느 경로로 답해도", "판정 수치는 완전히 동일합니다."]),
    ]
    gw = (CONTENT_W - 2 * 0.30) / 3
    for i, (t, lines) in enumerate(guards):
        x = CONTENT_X + i * (gw + 0.30)
        rect(slide, x, 5.42, gw, 1.00, fill=PANEL, line=LINE, radius=0.06,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.24, 0.14, 0.20, 0.08),
             items=[(t, 11, True, DARK)] +
                   [(l, 9.5, False, DARK_2, 6 if j == 0 else 3) for j, l in enumerate(lines)])

    txbox(slide, CONTENT_X, 6.50, CONTENT_W, 0.32,
          [("응답에는 reply · toolCalls[] · mode · provider가 함께 담겨, 어떤 프로바이더가 어떤 도구를 왜 호출했는지 검증할 수 있습니다.",
            10.5, True, GRAY)])
    return slide


# ==========================================================================
# 본선에서 지은 것 — 예선 덱에는 한 장도 없던 넷 (2026-08-16 신설)
#
# 이 덱은 15장 내내 **예선 프로토타입의 구조**를 설명하고 있었다. 본선에서 지은 것 —
# 정책 추출 파이프라인(2단계) · 규칙 관리 화면과 승인(4·5단계) · 감사·관측과 계보(7단계) ·
# 오프라인 재현(8단계) — 은 **한 글자도 없었다.** 전수 검색으로 확인했다:
# 「추출」 0 · 「승인」 0 · 「신고」 0 · 「계보」 0.
#
# 네 장의 공통 규율은 이 저장소의 규율과 같다 — **수를 문장에 박지 않고**, 실측하지
# 않은 것을 적지 않으며, 실패를 숨기지 않고 그 처리 방식을 보인다.
# ==========================================================================


def slide_13_extraction(prs):
    slide = new_slide(prs, 13, "정책 규칙 추출 파이프라인",
                      "공고문 원문에서 규칙을 뽑되, 근거가 원문에 없으면 그 초안을 통째로 버립니다.",
                      eyebrow="RULE EXTRACTION  ·  본선 확장")

    _engine_header(slide, "P1", "정책 규칙 추출", "Extraction",
                   "LLM 이 뽑고, 검증기가 원문과 대조하고, 사람이 승인해야 판정에 들어간다")

    # 줄을 상자 폭에 맞춰 짧게 유지한다 — 길면 줄바꿈이 생겨 4줄이 6줄이 되고,
    # `verify` 의 넘침 경고가 그 자리를 정확히 짚는다 (첫 판이 5건 걸렸다).
    _io_row(slide, 2.20, 1.55,
            ["공고문 원문 (policy_source)", "이용조건 · 출처표시",
             "정책 식별자 · 조회 시점", "원문은 코드포인트로 다룬다"],
            ["① LLM 이 초안 + 근거 span 제안", "② 검증기가 인용을 원문과 대조",
             "③ 하나라도 어긋나면 전체 폐기", "④ 통과분만 pending 으로 적재"],
            ["rule_draft (pending) + span", "실패 시 extraction_failed",
             "감사기록에 시도 · 지연 · 사유", "아직 판정에 들어가지 않는다"])

    # --- 이 파이프라인의 요지 -------------------------------------------
    txbox(slide, CONTENT_X, 3.88, 6.2, 0.30, [("이 파이프라인의 요지", 11, True, DARK)])
    rect(slide, CONTENT_X, 4.20, 5.86, 1.46, fill=RED_LT, line=RED, radius=0.06,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.26, 0.16, 0.22, 0.12),
         items=[("부분 저장을 하지 않는다", 12.5, True, RED),
                ("근거 인용 하나라도 원문에서 확인되지 않으면", 10.5, False, DARK_2, 6),
                ("그 초안 전체를 버립니다. 「절반만 맞는 규칙」을", 10.5, False, DARK_2, 3),
                ("만들지 않기 위해서입니다.", 10.5, False, DARK_2, 3)])

    rect(slide, CONTENT_X + 6.17, 4.20, 5.86, 1.46, fill=PANEL, line=LINE, radius=0.06,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.26, 0.16, 0.22, 0.12),
         items=[("실패는 숨기지 않고 이름을 붙인다", 12.5, True, DARK),
                ("span_not_in_text : 인용이 원문에 없다", 10, False, DARK_2, 6),
                ("span_missing : 값은 실었는데 근거가 없다", 10, False, DARK_2, 3),
                ("어긋난 위치(JSON 포인터)를 함께 남깁니다.", 10, False, GRAY, 3)])

    note_bar(slide, 5.80,
             "추출된 초안은 그 자체로는 아무 판정에도 참여하지 않습니다 — 규칙관리자가 "
             "승인해야 비로소 규칙 버전이 됩니다 (다음 장).")
    txbox(slide, CONTENT_X, 6.46, CONTENT_W, 0.30,
          [("span 오프셋의 단위는 유니코드 코드포인트이고, 원문을 자르는 일은 서버가 합니다.",
            10.5, False, GRAY)])
    return slide


def slide_14_admin_approval(prs):
    slide = new_slide(prs, 14, "규칙 관리 화면 · 승인 절차 · 권한",
                      "규칙을 바꾸는 유일한 경로는 사람의 승인이며, 그 권한은 서버가 강제합니다.",
                      eyebrow="REVIEW & APPROVAL  ·  본선 확장")

    # --- 3단 흐름 --------------------------------------------------------
    fw = (CONTENT_W - 2 * 0.52) / 3
    steps = [
        ("① 대기 큐", NAVY,
         ["추출 초안이 pending 으로 쌓입니다",
          "실패한 초안은 사유와 어긋난 위치를",
          "함께 보이고 승인 버튼이 닫힙니다",
          "현장 신고는 별도 유형으로 섭니다"]),
        ("② 검토 화면", YELLOW_DK,
         ["원문과 근거 인용을 나란히 대조",
          "변경 전 / 후 필드 diff",
          "승인 시 판정이 어떻게 바뀌는지",
          "회귀 사례로 미리 보여줍니다"]),
        ("③ 승인 · 반려", GREEN,
         ["승인 → 불변 RuleVersion 생성",
          "즉시 판정에 참여합니다",
          "반려에는 사유가 필수입니다",
          "둘 다 승인기록이 남습니다"]),
    ]
    for i, (title, accent, lines) in enumerate(steps):
        x = CONTENT_X + i * (fw + 0.52)
        rect(slide, x, 1.62, fw, 1.86, fill=WHITE, line=LINE, radius=0.06,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.28, 0.50, 0.22, 0.14),
             items=[(l, 10.5, False, DARK_2, 0 if j == 0 else 4)
                    for j, l in enumerate(lines)])
        badge(slide, x + 0.22, 1.76, 1.42, 0.30, title, fill=accent,
              color=WHITE if accent != YELLOW else DARK, size=9.5)
        if i < 2:
            arrow(slide, x + fw + 0.08, 2.40, 0.36, 0.30, "right", fill=YELLOW)

    # --- 권한 --------------------------------------------------------
    txbox(slide, CONTENT_X, 3.62, 6.0, 0.30, [("권한은 화면이 아니라 서버가 진다", 11, True, DARK)])
    roles = [
        ("상담원 (counselor)", ["판정 화면의 내부 정보를 봅니다",
                             "현장에서 발견한 이상을 신고합니다",
                             "규칙 승인 API 호출은 403 으로 거부됩니다"]),
        ("규칙관리자 (rule_manager)", ["대기 큐와 검토 화면에 들어갑니다",
                                  "승인·반려로 규칙 버전을 만듭니다",
                                  "관측 지표 화면을 함께 봅니다"]),
    ]
    rw = (CONTENT_W - 0.36) / 2
    for i, (t, lines) in enumerate(roles):
        rect(slide, CONTENT_X + i * (rw + 0.36), 3.94, rw, 1.24,
             fill=PANEL, line=LINE, radius=0.06,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.26, 0.16, 0.22, 0.12),
             items=[(t, 12, True, YELLOW_DK)] +
                   [(l, 10, False, DARK_2, 6 if j == 0 else 3)
                    for j, l in enumerate(lines)])

    rect(slide, CONTENT_X, 5.32, CONTENT_W, 0.62, fill=NAVY_LT, line=None, radius=0.08,
         align=PP_ALIGN.LEFT, insets=(0.26, 0.06, 0.24, 0.06),
         items=[("같은 초안을 둘이 동시에 승인해도 규칙 버전은 하나만 생깁니다 — "
                 "상태 검사와 상태 쓰기가 한 문장 안에 있는 조건부 갱신으로 선점합니다.",
                 10.5, True, NAVY)])
    txbox(slide, CONTENT_X, 6.10, CONTENT_W, 0.52,
          [("역할 분리는 화면을 나눠서 얻는 것이 아닙니다. 화면은 하나이고, 막는 것은 서버입니다 — "
            "상담원 세션으로 승인 API 를 직접 호출하면 403 입니다.", 10.5, False, GRAY)])
    return slide


def slide_15_audit_provenance(prs):
    slide = new_slide(prs, 15, "감사 추적 · 관측 지표 · 값의 계보",
                      "무엇이 이 숫자를 말했는가 — 값마다 출처와 검증 상태가 함께 다닙니다.",
                      eyebrow="AUDIT & PROVENANCE  ·  본선 확장")

    # --- 계보 3분류 ------------------------------------------------------
    txbox(slide, CONTENT_X, 1.58, 6.4, 0.30, [("값의 계보 (Provenance)", 11.5, True, DARK)])
    grades = [
        ("verified", GREEN, GREEN_LT,
         ["출처 문서를 열어 대조한 값", "예) 지역 시세 5필드 —", "국토교통부 실거래가 실측"]),
        ("unverified", AMBER, YELLOW_LT,
         ["출처를 아직 특정하지 못한 값", "예) 관리비 · 시장 위험도 ·", "정책 조건과 금리의 시드값"]),
        ("our_choice", NAVY, NAVY_LT,
         ["공표 준거가 없어 우리가 정한 값", "예) 판정 임계·가중치 상당수", "그 사실을 그대로 적습니다"]),
    ]
    gw = (CONTENT_W - 2 * 0.30) / 3
    for i, (name, accent, bg, lines) in enumerate(grades):
        x = CONTENT_X + i * (gw + 0.30)
        rect(slide, x, 1.92, gw, 1.34, fill=bg, line=None, radius=0.06,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.26, 0.16, 0.22, 0.12),
             items=[(name, 12.5, True, accent)] +
                   [(l, 10, False, DARK_2, 6 if j == 0 else 3)
                    for j, l in enumerate(lines)])

    txbox(slide, CONTENT_X, 3.36, CONTENT_W, 0.30,
          [("화면은 이 계보를 등급과 함께 보이고, 등급이 낮은 이유를 원인 유형별로 나눠 적습니다 — "
            "「검증 안 됨 N건, 대응 주체는 수집 배치」처럼.", 10.5, False, GRAY)])

    # --- 감사 · 관측 -----------------------------------------------------
    panes = [
        ("감사 추적 (AuditEvent)",
         ["추출 · 승인 · 반려 · 신고가 남습니다",
          "덧붙이기만 가능하며 지우는 경로가 없습니다",
          "지운 기록은 감사 추적이 아니기 때문입니다",
          "사유 본문은 파일 로그에 찍지 않습니다"]),
        ("관측 지표",
         ["배치 상태 · 데이터 신선도 · 대기 큐",
          "LLM 호출 성공률과 지연",
          "★ 없는 값을 0 으로 그리지 않습니다",
          "임계가 미정인 것에는 판정선을 긋지 않습니다"]),
    ]
    pw = (CONTENT_W - 0.36) / 2
    for i, (t, lines) in enumerate(panes):
        rect(slide, CONTENT_X + i * (pw + 0.36), 3.76, pw, 1.56,
             fill=WHITE, line=LINE, radius=0.06,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.28, 0.18, 0.22, 0.12),
             items=[(t, 12, True, DARK)] +
                   [(l, 10, False, DARK_2, 6 if j == 0 else 3)
                    for j, l in enumerate(lines)])

    rect(slide, CONTENT_X, 5.48, CONTENT_W, 0.72, fill=DARK, line=None, radius=0.08,
         align=PP_ALIGN.LEFT, insets=(0.30, 0.08, 0.26, 0.08),
         items=[("값을 지어내지 않습니다. 빈 칸은 실패가 아닙니다.", 13, True, YELLOW),
                ("출처를 못 찾은 항목은 못 찾았다고 적고, 준거가 없는 선택은 우리 선택이라고 적습니다.",
                 10.5, False, WHITE, 4)])
    txbox(slide, CONTENT_X, 6.34, CONTENT_W, 0.30,
          [("이 원칙이 화면·API 응답·기술설명서에 같은 문장으로 나타나는지를 자동 검사가 대조합니다.",
            10.5, False, GRAY)])
    return slide


def slide_16_offline_rehearsal(prs):
    slide = new_slide(prs, 16, "오프라인 동작과 시연 재현",
                      "네트워크를 실제로 끊고 같은 장면을 3회 재현했습니다.",
                      eyebrow="OFFLINE & REPRODUCIBILITY  ·  본선 확장")

    # --- 무엇을 어떻게 쟀나 ---------------------------------------------
    cols = [
        ("무엇을 끊었나", ["무선을 끄고 그 상태를 유지했습니다",
                      "회차마다 5개 지점에서 차단을 확인",
                      "판정은 HTTPS 요청 한 갈래로만"]),
        ("무엇이 같았나", ["승인 → 시민 화면 변화 장면",
                      "회차 간 대조 항목 32개",
                      "불일치 0건 · 장면 오류 0"]),
        ("무엇이 실패해도 되나", ["네트워크가 필요한 것은 수집 배치뿐",
                          "실패해도 이전 값이 유지되고",
                          "계보가 stale 로 내려갑니다"]),
    ]
    cw = (CONTENT_W - 2 * 0.30) / 3
    for i, (t, lines) in enumerate(cols):
        x = CONTENT_X + i * (cw + 0.30)
        rect(slide, x, 1.62, cw, 1.42, fill=PANEL, line=LINE, radius=0.06,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.26, 0.16, 0.22, 0.12),
             items=[(t, 12, True, YELLOW_DK)] +
                   [(l, 10, False, DARK_2, 6 if j == 0 else 3)
                    for j, l in enumerate(lines)])

    # --- 이 리허설의 값 --------------------------------------------------
    rect(slide, CONTENT_X, 3.20, CONTENT_W, 1.24, fill=RED_LT, line=RED, radius=0.06,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.30, 0.16, 0.26, 0.12),
         items=[("첫 판정을 우리 손으로 뒤집었습니다", 12.5, True, RED),
                ("처음에는 「3회 무오류」로 끝냈습니다. 그런데 차단을 확인해 보니 "
                 "무선이 스스로 되살아나 있었습니다 — 「끊고 시작한 3회」였지 "
                 "「끊은 상태의 3회」가 아니었습니다.", 10.5, False, DARK_2, 6),
                ("그래서 차단 방식을 바꾸고 회차마다 차단을 다시 확인한 뒤 처음부터 "
                 "다시 돌았습니다. 통과한 회차를 골라 굳히지 않았습니다.",
                 10.5, False, DARK_2, 3)])

    # --- 시연 구성 -------------------------------------------------------
    txbox(slide, CONTENT_X, 4.62, 6.0, 0.30, [("시연은 노트북 한 대에서 끝납니다", 11, True, DARK)])
    stack = [
        ("준비", ["저장소를 픽스처로 굳혀 커밋했습니다",
                "클론 후 실행만으로 같은 상태가 섭니다"]),
        ("무대", ["시민 판정 → 상담원 신고 → 규칙 승인",
                "→ 시민 화면 변화까지 한 흐름"]),
        # ★ 「증거 파일을 남겼습니다」라고 적으려다 확인하고 고쳤다 — 그 파일들은
        #   `backend/var/` 아래의 런타임 산출물이고 **저장소에 없다**(.gitignore).
        #   저장소에 실제로 남아 있는 것은 **대본 문서의 관측 기록**이다.
        #   확인할 수 없는 것을 덱에 적지 않는다.
        ("기록", ["회차별 관측을 대본 문서에 적었습니다",
                "채택하지 않은 회차와 그 사유까지"]),
    ]
    sw2 = (CONTENT_W - 2 * 0.30) / 3
    for i, (t, lines) in enumerate(stack):
        x = CONTENT_X + i * (sw2 + 0.30)
        rect(slide, x, 4.96, sw2, 1.02, fill=WHITE, line=LINE, radius=0.06,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.26, 0.14, 0.22, 0.10),
             items=[(t, 11.5, True, DARK)] +
                   [(l, 10, False, DARK_2, 5 if j == 0 else 3)
                    for j, l in enumerate(lines)])

    note_bar(slide, 6.14,
             "재현이 값을 낸 자리는 통과가 아니라 실패였습니다 — 대본이 화면을 보지 않고 "
             "적어 둔 문장 둘을 이 재현이 걷어냈습니다.")
    return slide


def slide_17_screens_stack(prs):
    slide = new_slide(prs, 17, "화면 흐름과 기술 스택",
                      "단일 페이지 3단 구성 — 입력에서 상담까지 이동 없이 이어집니다.",
                      eyebrow="UX FLOW & TECH STACK")

    sw = (CONTENT_W - 2 * 0.24) / 3
    screens = [
        ("① 프로필 입력", ["나이 · 소득 · 자산 · 부채", "가구원수 · 희망 지역", "무주택 / 신혼 / 중소기업 재직"],
         [(0.55, 0.14), (0.55, 0.14), (0.55, 0.14), (0.30, 0.14)]),
        ("② 결과 대시보드", ["시나리오별 부담 막대 (인라인 SVG)", "시나리오 비교 카드 · 적합도",
                          "정책 매칭 리스트 · 리스크 스코어"],
         [(0.55, 0.22), (0.26, 0.16), (0.26, 0.16), (0.55, 0.14)]),
        ("③ AI 상담 채팅", ["자연어 후속 질문", "엔진 결과 기반 답변", "호출된 도구 표시"],
         [(0.55, 0.12), (0.38, 0.12), (0.55, 0.12), (0.55, 0.16)]),
    ]
    for i, (title, caps, rows) in enumerate(screens):
        x = CONTENT_X + i * (sw + 0.24)
        rect(slide, x, 1.58, sw, 2.62, fill=WHITE, line=LINE, radius=0.05)
        rect(slide, x, 1.58, sw, 0.42, fill=DARK, line=None, radius=0.05,
             items=[(title, 11.5, True, WHITE)], insets=(0.08, 0.02, 0.08, 0.02))
        # wireframe blocks
        wy = 2.14
        inner_x = x + 0.22
        inner_w = sw - 0.44
        for frac, hh in rows:
            rect(slide, inner_x, wy, inner_w * frac, hh, fill=PANEL, line=None,
                 radius=0.20)
            wy += hh + 0.10
        rect(slide, inner_x, wy + 0.02, inner_w * 0.42, 0.20, fill=YELLOW, line=None,
             radius=0.30)
        cy = 3.42
        txbox(slide, inner_x, cy, inner_w, 0.72,
              [(c, 9.5, False, DARK_2, 0 if j == 0 else 3) for j, c in enumerate(caps)])
        if i < 2:
            arrow(slide, x + sw + 0.02, 2.72, 0.20, 0.26, "right", fill=YELLOW)

    # 「내장 MOCK 응답」은 없다. 백엔드가 없을 때 도는 것은 생성물 기반 로컬 엔진이고,
    # 그 사실을 화면이 배너로 밝힌다 (SPEC 6.2 D-11 #3 · PR #67 ⑥).
    txbox(slide, CONTENT_X, 4.30, CONTENT_W, 0.30,
          [("스크롤 한 번으로 입력 → 판정 → 상담이 이어지며, 백엔드 미기동 시에도 생성물 기반 로컬 판정으로 화면 시연이 가능합니다.",
            10.5, False, GRAY)])

    stacks = [
        ("프론트엔드", ["Vanilla HTML / CSS / JS", "인라인 SVG 차트 직접 생성",
                     "빌드툴 · CDN 미사용 (오프라인 동작)", "반응형 · 한국어 UI"]),
        # 선언된 하한은 3.11 이고(backend/pyproject.toml), 판정이 읽는 데이터는
        # `data/*.json` 이 아니라 저장소다. 실행 스크립트는 scripts/dev.bat 이다.
        ("백엔드", ["Python 3.11+ · FastAPI", "engines/ : 순수 함수 4종",
                  "store/ : SQLite (data/ 는 시드)", "uvicorn + dev.bat 원클릭 실행"]),
        ("AI 레이어", ["LLM 프로바이더 추상화 계층", "1순위 OpenAI function calling",
                    "2순위 Anthropic tool use", "3순위 룰 기반 결정론 폴백"]),
        # ★ 「pytest 단위 테스트」는 거짓은 아니나 **과소**였다. 실제로는 교차 테스트가
        #   계약·아키텍처·생성물 바이트·오프라인 기동을 함께 붙들고, 네트워크를 끊은
        #   3회 재현이 따로 있다. 수는 적지 않는다 — 적으면 낡는다(이 덱이 이미 겪었다).
        ("품질 · 검증", ["pytest 단위·교차 테스트 (경계값 포함)", "동일 입력 → 동일 출력 재현성",
                      "계약·아키텍처·생성물 자동 대조", "네트워크 차단 상태 3회 재현"]),
    ]
    tw = (CONTENT_W - 3 * 0.28) / 4
    for i, (t, lines) in enumerate(stacks):
        x = CONTENT_X + i * (tw + 0.28)
        rect(slide, x, 4.72, tw, 1.52, fill=PANEL, line=LINE, radius=0.06,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.22, 0.16, 0.18, 0.10),
             items=[(t, 12, True, YELLOW_DK)] +
                   [(l, 9.5, False, DARK_2, 6 if j == 0 else 3) for j, l in enumerate(lines)])
    txbox(slide, CONTENT_X, 6.38, CONTENT_W, 0.30,
          [("외부 네트워크 의존을 최소화해 심사 환경에서 재현 가능한 실행을 우선했습니다.", 10.5, True, GRAY)])
    return slide


def slide_18_data_notice(prs):
    slide = new_slide(prs, 18, "데이터 정직성 고지",
                      "프로토타입에 사용된 수치의 성격과 한계를 명확히 밝힙니다.",
                      eyebrow="DATA DISCLOSURE")

    rect(slide, CONTENT_X, 1.58, CONTENT_W, 0.86, fill=RED_LT, line=RED, line_w=1.2,
         radius=0.08, align=PP_ALIGN.LEFT, insets=(0.30, 0.10, 0.26, 0.10),
         items=[("필수 고지", 10, True, RED),
                (DISCLAIMER, 14, True, DARK, 5)])

    items = [
        ("금리 · 한도 수치", ["금리·한도·기간 값은 모두", "시연을 위한 예시 수치입니다.",
                          "특정 금융회사의 실제 조건을", "옮기거나 추정하지 않았습니다."]),
        ("제도 요건 기술", ["연령(만 19~34세), 무주택 등", "공개된 일반 요건 수준으로만",
                        "기술했습니다. 세부 심사 기준은", "취급 기관·시점에 따라 다릅니다."]),
        ("통계 인용 원칙", ["출처를 명시할 수 없는 통계", "수치는 인용하지 않았습니다.",
                        "문제 정의는 정성적 서술로", "대체하고 참고 자료만 표기."]),
        # ★ 「구조 시연용 예시 데이터」는 **이제 거짓이다.** 계약 결정 #40 이 시세를
        #   실수집으로 굳혀 8필드 중 5개(전세중위가·월세보증금·월세·전세가율·전환율)가
        #   국토교통부 실거래가에서 온 `verified` 다. 시민 화면 푸터는 이미 이렇게
        #   말하고 있었고 **덱만 낡아 있었다** — 리허설이 화면에서 잡은 것과 같은 부류다.
        ("지역 시세 데이터", ["항목마다 다릅니다. 8필드 중 5개는", "국토교통부 실거래가 실측값이고,",
                         "관리비·시장위험도·보증가능은", "출처를 특정하지 못한 값입니다."]),
    ]
    cw = (CONTENT_W - 3 * 0.28) / 4
    for i, (t, lines) in enumerate(items):
        x = CONTENT_X + i * (cw + 0.28)
        rect(slide, x, 2.60, cw, 1.86, fill=WHITE, line=LINE, radius=0.06,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.24, 0.18, 0.20, 0.12),
             items=[(t, 12, True, DARK)] +
                   [(l, 9.5, False, DARK_2, 8 if j == 0 else 3) for j, l in enumerate(lines)])
        rect(slide, x, 2.60, cw, 0.08, fill=RED, line=None, shape_type=MSO_SHAPE.RECTANGLE)

    # `disclaimer` 를 **항목마다** 갖는 것은 policies.json 이다. regions.json 은 항목에
    # `source` 만 있고 고지는 파일 최상위 `_disclaimer` 한 벌이다. 「모든 항목」이라고
    # 적으면 그 절반이 거짓이 된다.
    rect(slide, CONTENT_X, 4.72, 7.60, 1.52, fill=PANEL, line=LINE, radius=0.06,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.28, 0.18, 0.24, 0.12),
         items=[("데이터 구조로 강제한 고지 — data/*.json 필수 필드", 12, True, DARK),
                ("\"source\"       : 출처 기관명 (예: \"주택도시기금\", \"국토교통부\") — 전 항목", 11, False, DARK_2, 7),
                ("\"disclaimer\" : 정책 항목마다 동일 포함 (시세는 파일 단위 고지)", 11, False, DARK_2, 4),
                ("→ 데이터에 출처와 고지가 없으면 화면에도 노출되지 않는 구조입니다.",
                 11, True, YELLOW_DK, 6)])
    rect(slide, 8.53, 4.72, 4.15, 1.52, fill=DARK, line=None, radius=0.06,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.26, 0.18, 0.22, 0.12),
         items=[("고지 노출 위치", 11.5, True, YELLOW),
                ("· 서비스 화면 하단 상시 노출", 10.5, False, WHITE, 7),
                ("· API 응답 meta.disclaimer", 10.5, False, WHITE, 4),
                ("· 정책 항목별 disclaimer 필드", 10.5, False, WHITE, 4),
                ("· 본 기술설명서 18장", 10.5, False, WHITE, 4)])
    txbox(slide, CONTENT_X, 6.40, CONTENT_W, 0.30,
          [("본 산출물은 의사결정을 보조하는 참고 도구이며, 금융상품 청약·계약 권유나 투자 자문이 아닙니다.",
            10.5, True, GRAY)])
    return slide


def slide_19_impact(prs):
    slide = new_slide(prs, 19, "금융 서비스 연계 · 기대효과 · 향후 확장",
                      "주택금융 접점에서 바로 쓰일 수 있는 판정 엔진을 지향합니다.",
                      eyebrow="BUSINESS FIT & ROADMAP")

    links = [
        ("사업 연결성", ["은행·보증기관의 핵심 사업영역인", "주택금융·전세자금대출 접점과 직결됩니다.",
                     "첫 주거를 준비하는 청년 고객의", "첫 접점 도구로 활용할 수 있습니다."]),
        ("상담 품질", ["판정 결과와 근거 문자열이 함께 남아", "상담원이 설명 근거로 그대로 사용하고",
                    "고객은 상담 전에 자기 상황을", "정량적으로 파악한 채 방문합니다."]),
        ("도입 유연성", ["모든 판정에 rationale이 동반되어", "동일 입력에 동일 결과가 재현되고,",
                    "LLM은 프로바이더 교체만으로", "내부망 전용 모델로 전환 가능합니다."]),
    ]
    cw = (CONTENT_W - 2 * 0.30) / 3
    for i, (t, lines) in enumerate(links):
        x = CONTENT_X + i * (cw + 0.30)
        rect(slide, x, 1.58, cw, 1.86, fill=PANEL, line=LINE, radius=0.06,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.26, 0.44, 0.22, 0.12),
             items=[(t, 12.5, True, DARK)] +
                   [(l, 10.5, False, DARK_2, 7 if j == 0 else 2) for j, l in enumerate(lines)])
        rect(slide, x, 1.58, 0.075, 1.86, fill=YELLOW, line=None,
             shape_type=MSO_SHAPE.RECTANGLE)
        badge(slide, x + 0.26, 1.74, 0.44, 0.24, "0%d" % (i + 1), fill=DARK, color=YELLOW,
              size=8.5)

    txbox(slide, CONTENT_X, 3.62, 6.0, 0.30, [("향후 확장 로드맵", 12.5, True, DARK)])
    steps = [("Phase 1", "프로토타입", ["4대 엔진 · 대시보드", "· AI 상담 (현재)"]),
             ("Phase 2", "데이터 연동", ["공신력 있는 시세·정책", "데이터 소스 연결"]),
             ("Phase 3", "상품 연계", ["기관별 주택금융 상품 카탈로그", "· 내부망 모델 전환"]),
             ("Phase 4", "사후 관리", ["계약 이후 주거비 모니터링", "· 갱신 시점 재판정"])]
    pw = (CONTENT_W - 3 * 0.40) / 4
    for i, (ph, t, lines) in enumerate(steps):
        x = CONTENT_X + i * (pw + 0.40)
        fill = YELLOW_LT if i == 0 else WHITE
        rect(slide, x, 3.98, pw, 1.24, fill=fill, line=LINE, radius=0.06,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.22, 0.40, 0.18, 0.10),
             items=[(t, 12, True, DARK)] +
                   [(l, 10, False, DARK_2, 6 if j == 0 else 2) for j, l in enumerate(lines)])
        badge(slide, x + 0.18, 4.12, 0.86, 0.26, ph, fill=DARK, color=YELLOW, size=8.5)
        if i < 3:
            arrow(slide, x + pw + 0.04, 4.46, 0.32, 0.28, "right", fill=YELLOW)

    rect(slide, CONTENT_X, 5.44, CONTENT_W, 1.10, fill=DARK, line=None, radius=0.08,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, insets=(0.34, 0.18, 0.30, 0.12),
         items=[("맺음말", 10, True, YELLOW),
                ("청년에게 필요한 것은 더 많은 상품 정보가 아니라, “나는 지금 얼마짜리 집에 살아도 되는가”에 대한 답입니다.",
                 14, True, WHITE, 6),
                ("Home_Compass는 그 답을 숫자와 근거로 제시하고, 실행 가능한 다음 단계까지 연결합니다.",
                 12, False, GRAY_LT, 5)])
    return slide


# ==========================================================================
# build & verify
# ==========================================================================


def _endpoint_counts() -> tuple[int, int, int]:
    """계약(`contracts/openapi.json`)에서 엔드포인트를 **센다.** 손으로 적지 않는다.

    PR #76 이 「15종」이라고 적어 둔 뒤 6-A 신고 둘(`POST /api/reports` ·
    `GET /api/admin/reports`)과 7단계 감사(`GET /api/admin/audit`)가 늘어 **18종이
    됐는데 아무도 못 잡았다.** 수를 문장에 박으면 낡는다 — 이 저장소가 같은 부류를
    반복해서 잡아왔고(`ingest` 주석의 10건 · `regions.js` 의 시점 주석) 이 덱도 그랬다.

    돌려주는 것은 (시민 판정 경로 수, 규칙관리 수, 전체 수).
    """
    with open(os.path.join(_REPO_ROOT, "contracts", "openapi.json"), encoding="utf-8") as fh:
        spec = json.load(fh)
    methods = {"get", "post", "put", "delete", "patch"}
    ops = [(p, m) for p, item in spec["paths"].items() for m in item if m in methods]
    admin = sum(1 for p, _ in ops if p.startswith("/api/admin/"))
    auth = sum(1 for p, _ in ops if p.startswith("/api/auth/"))
    return len(ops) - admin - auth, admin, len(ops)


def build() -> str:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    cp = prs.core_properties
    cp.title = "Home_Compass 기술설명서"
    cp.subject = "2026 금융 AI Challenge"
    cp.author = "Home_Compass"

    builders = [
        slide_01_cover, slide_02_problem, slide_03_gap, slide_04_overview,
        slide_05_diff, slide_06_architecture, slide_07_pipeline, slide_08_e1,
        slide_09_e2, slide_10_e3, slide_11_e4, slide_12_llm,
        # 본선에서 지은 것 넷. 예선 덱에는 한 장도 없었다 (2026-08-16 신설).
        slide_13_extraction, slide_14_admin_approval,
        slide_15_audit_provenance, slide_16_offline_rehearsal,
        slide_17_screens_stack, slide_18_data_notice, slide_19_impact,
    ]
    # ★ 장수를 손으로 적지 않는다 — 이 목록이 정본이고 쪽번호가 그것을 받아 쓴다.
    global TOTAL_SLIDES
    TOTAL_SLIDES = len(builders)
    for fn in builders:
        fn(prs)

    prs.save(OUT_PATH)
    return OUT_PATH


# --- verification ---------------------------------------------------------

_EMU_PER_IN = 914400


def _shape_text_stats(shape):
    """Return (chars, max_font_pt, paragraph_count) for a shape's text frame."""
    if not shape.has_text_frame:
        return 0, 0, 0
    chars = 0
    max_pt = 0.0
    paras = 0
    for p in shape.text_frame.paragraphs:
        line_chars = 0
        for r in p.runs:
            t = r.text or ""
            line_chars += len(t)
            if r.font.size is not None:
                max_pt = max(max_pt, r.font.size.pt)
        if line_chars > 0:
            paras += 1
            chars += line_chars
    return chars, max_pt, paras


# Line pitch factor measured from PowerPoint renders of this deck:
# a 10.5pt run with 4pt space_before occupied ~18pt, i.e. 1.33 * size + space.
_LINE_PITCH = 1.33
# Shapes do not clip text, so a few points of spill is invisible. Only flag
# boxes that exceed usable height by more than this ratio -- calibrated so the
# flags match what was actually visible in the rendered PNGs.
_OVERFLOW_TOLERANCE = 1.15


def _estimate_overflow(shape):
    """Advisory overflow heuristic: does the text plausibly fit the box height?

    Hangul glyphs are ~1.0em wide, latin ~0.52em. This is an estimate, not a
    renderer -- treat hits as "worth eyeballing", not as proof of breakage.
    """
    if not shape.has_text_frame:
        return None
    tf = shape.text_frame
    w_in = shape.width / _EMU_PER_IN
    h_in = shape.height / _EMU_PER_IN
    ml = (tf.margin_left or 0) / _EMU_PER_IN
    mr = (tf.margin_right or 0) / _EMU_PER_IN
    mt = (tf.margin_top or 0) / _EMU_PER_IN
    mb = (tf.margin_bottom or 0) / _EMU_PER_IN
    usable_w_pt = max((w_in - ml - mr) * 72.0, 1.0)
    usable_h_pt = max((h_in - mt - mb) * 72.0, 1.0)

    total_h = 0.0
    for p in tf.paragraphs:
        text = "".join(r.text or "" for r in p.runs)
        if not text.strip():
            continue
        size = 12.0
        for r in p.runs:
            if r.font.size is not None:
                size = max(size, r.font.size.pt)
        em = 0.0
        for ch in text:
            em += 1.0 if ord(ch) > 0x1100 else 0.52
        width_pt = em * size
        lines = max(1, math.ceil(width_pt / usable_w_pt))
        sb = p.space_before.pt if p.space_before is not None else 0.0
        total_h += lines * size * _LINE_PITCH + sb
    if total_h == 0:
        return None
    return total_h, usable_h_pt


def verify(path: str = OUT_PATH) -> int:
    prs = Presentation(path)
    slides = list(prs.slides)
    print("=" * 78)
    print("VERIFY : %s" % path)
    print("file size      : {:,} bytes".format(os.path.getsize(path)))
    print("slide size     : %.3f x %.3f in  (ratio %.4f)"
          % (prs.slide_width / _EMU_PER_IN, prs.slide_height / _EMU_PER_IN,
             prs.slide_width / prs.slide_height))
    print("slide count    : %d" % len(slides))
    print("=" * 78)

    problems = []   # hard failures -> non-zero exit
    advisories = []  # estimated layout hints -> printed only
    # ★ 「12~15」라는 범위가 박혀 있었다. 본선 산출물 넷을 더하자 그 자리가 빨간불이 됐고,
    #   **그것이 옳게 작동한 것이다** — 다만 고칠 방향은 범위를 넓히는 것이 아니라
    #   생성기가 실제로 만드는 수와 대조하는 것이다. 수를 문장에 박으면 낡는다.
    #   (`--verify` 로 단독 실행하면 `TOTAL_SLIDES` 가 0 이라 대조할 정본이 없다.
    #    그때는 이 항목을 건너뛰는 대신 **그 사실을 문제로 적는다** — 침묵하지 않는다.)
    if TOTAL_SLIDES == 0:
        problems.append("빌드 없이 --verify 만 돌려 장수의 정본을 모른다 "
                        "(python build_ppt.py 로 함께 돌린다)")
    elif len(slides) != TOTAL_SLIDES:
        problems.append("slide count %d != builders %d" % (len(slides), TOTAL_SLIDES))
    if abs(prs.slide_width / prs.slide_height - 16 / 9) > 0.01:
        problems.append("aspect ratio is not 16:9")

    bad_font_runs = 0
    total_runs = 0
    total_shapes = 0
    total_chars = 0

    for i, slide in enumerate(slides, start=1):
        shapes = list(slide.shapes)
        total_shapes += len(shapes)
        drawn = sum(1 for s in shapes if s.shape_type is not None
                    and not s.has_text_frame or s.shape_type is not None)
        texts = []
        n_auto = 0
        for s in shapes:
            if not s.has_text_frame:
                continue
            for p in s.text_frame.paragraphs:
                for r in p.runs:
                    total_runs += 1
                    t = r.text or ""
                    total_chars += len(t)
                    if r.font.name != FONT:
                        bad_font_runs += 1
                    if t.strip():
                        texts.append(t.strip())
            ov = _estimate_overflow(s)
            if ov and ov[0] > ov[1] * _OVERFLOW_TOLERANCE:
                n_auto += 1
                advisories.append(
                    "slide %02d: check text fit (est. ~%.0fpt / box %.0fpt) : %s"
                    % (i, ov[0], ov[1], texts[-1][:38] if texts else "")
                )
        n_shapes_drawn = sum(
            1 for s in shapes
            if s.shape_type is not None and str(s.shape_type).startswith("AUTO_SHAPE")
        )
        n_conn = sum(1 for s in shapes if s.shape_type is not None
                     and "LINE" in str(s.shape_type).upper())
        print("\n--- slide %02d --- shapes=%d (autoshape=%d, connector=%d) chars=%d"
              % (i, len(shapes), n_shapes_drawn, n_conn,
                 sum(len(t) for t in texts)))
        for t in texts:
            print("    | %s" % t)

    print("\n" + "=" * 78)
    print("total shapes   : %d" % total_shapes)
    print("total runs     : %d   (chars=%d)" % (total_runs, total_chars))
    print("runs not in '%s' : %d" % (FONT, bad_font_runs))
    if bad_font_runs:
        problems.append("%d runs are not set to %s" % (bad_font_runs, FONT))

    if advisories:
        print("\nLAYOUT ADVISORIES (%d) — estimates only, confirm in a render:"
              % len(advisories))
        for a in advisories:
            print("  ~ %s" % a)

    if problems:
        print("\nPROBLEMS (%d):" % len(problems))
        for p in problems:
            print("  ! %s" % p)
    else:
        print("\nOK — no hard problems detected (slide count, 16:9, font coverage).")
    print("=" * 78)
    return len(problems)


def deck_text(path: str = OUT_PATH) -> list[str]:
    """덱의 본문을 순서대로 뜬다. **재생성 대조의 기준이 이것이다.**

    ★ 바이트로 비교할 수 없다 — 재 봤다. 같은 코드로 두 번 뽑아도 sha 가 다르다
      (`.pptx` 는 zip 이고 항목마다 시각이 박힌다). **본문 텍스트는 결정적이다** —
      두 번 뽑아 544문장이 완전히 일치했다. 그래서 `frontend/generated/` 의 바이트
      비교와 **같은 목적을 다른 수단으로** 이룬다.
    """
    prs = Presentation(path)
    out = []
    for index, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs)
                if text.strip():
                    out.append("%02d|%s" % (index, text))
    return out


def check() -> int:
    """커밋된 덱이 지금 코드로 다시 뽑은 것과 같은가. 다르면 1.

    **PR #76 이 정확히 여기서 샜다** — 생성기 22곳을 고치고 덱을 다시 뽑지 않아서
    커밋된 제출물에 옛 거짓이 그대로 남았고, 그것을 잡는 검사가 없었다.
    """
    if not os.path.exists(OUT_PATH):
        print("커밋된 덱이 없다: %s" % OUT_PATH)
        return 1
    committed = deck_text(OUT_PATH)
    import shutil
    import tempfile
    backup = None
    try:
        with tempfile.TemporaryDirectory() as tmp:
            backup = os.path.join(tmp, "committed.pptx")
            shutil.copy2(OUT_PATH, backup)
            build()
            fresh = deck_text(OUT_PATH)
            shutil.copy2(backup, OUT_PATH)   # 커밋본을 되돌린다 — `--check` 는 쓰지 않는다
    except Exception as exc:              # noqa: BLE001 — 무엇이 터지든 커밋본을 지키고 알린다
        print("재생성에 실패했다: %r" % (exc,))
        return 1
    if committed == fresh:
        print("OK  %s 는 build_ppt.py 와 일치한다 (%d문장)" % (OUT_PATH, len(fresh)))
        return 0
    only_committed = [t for t in committed if t not in set(fresh)]
    only_fresh = [t for t in fresh if t not in set(committed)]
    print("커밋된 덱이 코드와 다르다. `python build_ppt.py` 로 다시 뽑아 함께 커밋한다.")
    print("  커밋본에만 (%d):" % len(only_committed))
    for t in only_committed[:20]:
        print("    - %s" % t)
    print("  재생성에만 (%d):" % len(only_fresh))
    for t in only_fresh[:20]:
        print("    + %s" % t)
    return 1


if __name__ == "__main__":
    force_utf8_stdout()
    if "--check" in sys.argv:
        sys.exit(check())
    if "--verify" in sys.argv:
        sys.exit(0 if verify() == 0 else 1)
    out = build()
    print("built: %s" % out)
    # ★ 예전에는 `verify(out)` 의 결과를 받아 놓고 **버렸다** (`sys.exit(0)`).
    #   자체점검이 문제를 찍어도 종료코드가 0 이라 **무조건 통과하는 검사**였다 —
    #   이 저장소가 이름 붙인 실패 양상이고 여기서도 그대로였다.
    sys.exit(0 if verify(out) == 0 else 1)
